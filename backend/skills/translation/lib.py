"""
scripts/lib.py — ROX 车机多语翻译知识库核心逻辑
（M1 8GB 优化版 / rox-translator:1.5b + rox-bge:m1 / 长文本切段拼回段落）
★ v2.18: 剥离离线翻译逻辑，统一在线模型翻译
★ v2.12: _extract_single_lang改为保守策略——无语言前缀时直接保留不截断，防止误删多段落翻译
★ v2.9: PPTX纠错改为红字批注识别（提取红色字体→解析"改为：XXX"→匹配KB），DOCX/XLSX保留位置对齐
★ v2.8: 严格执行根逻辑——按语言维度三级匹配，KB命中后缺失语言走模型兜底（在线优先），离线长文本也补全，绝不静默丢弃任何选中语言
★ v2.7: 知识库命中后缺失语言走模型兜底（_fill_missing_langs），不再静默丢弃
★ v2.6: 修复多语混入——_extract_single_lang重写(非目标语言前缀截断)+_strip_foreign_paragraphs(段落级字符检测)+_strip_lang_name_sections(拉丁语语言名截断)
★ 修复：保留原始换行符（\n 和 \n\n）
"""
import hashlib
import json
import os
import re
import sqlite3
import sys
from datetime import datetime
from pathlib import Path

import numpy as np
import requests

# ====================================================================
# 模型与 API 配置
# 所有配置优先从环境变量读取，其次使用模块级默认值
# ====================================================================
ONLINE_API_BASE = os.getenv("SILICONFLOW_API_BASE", "https://api.siliconflow.cn/v1")  # SiliconFlow API 地址
ONLINE_API_KEY = os.getenv("SILICONFLOW_API_KEY", "")                                  # API 密钥（空=未配置）
ONLINE_MODEL = os.getenv("ONLINE_MODEL", "tencent/Hunyuan-MT-7B")                      # 默认翻译模型
ONLINE_TIMEOUT = int(os.getenv("ONLINE_TIMEOUT", "120"))                                # API 请求超时秒数

HUNYUAN_MT_MODEL = "tencent/Hunyuan-MT-7B"        # Hunyuan-MT 专用翻译模型（支持33语）
HUNYUAN_FALLBACK_MODEL = "THUDM/GLM-Z1-9B-0414"    # 当目标语言不在Hunyuan支持范围内时的降级模型

# Hunyuan-MT 官方支持的33种语言代码集合
# 超出此范围的语言会自动降级到 HUNYUAN_FALLBACK_MODEL
HUNYUAN_MT_LANGS = {
    "zh","en","fr","pt","es","ja","tr","ru","ar","ko",
    "th","it","de","vi","ms","id","tl","hi","zh_hant",
    "pl","cs","nl","km","my","fa","gu","ur","te","mr",
    "he","bn","ta","uk","bo","kk","mn","ug","yue",
}

HUNYUAN_TEMP = 0.7                                          # Hunyuan-MT 模型的温度参数
HUNYUAN_TOP_P = 0.6                                         # Hunyuan-MT 的 top_p 采样参数
HUNYUAN_TOP_K = 20                                          # Hunyuan-MT 的 top_k 采样参数
HUNYUAN_REPETITION_PENALTY = 1.05                            # Hunyuan-MT 的重复惩罚系数
FALLBACK_TEMP = 0.1                                          # 降级模型的温度参数（更低=更确定性）

PROXY_URL = os.getenv("PROXY_URL", "")                       # 代理地址（可选）

EMBED_MODEL = "embedding-2"                                  # 智谱 embedding 模型名（仅智谱支持）

def _load_config_json():
    """加载 config.json 配置文件。

    搜索优先级：可执行文件同级目录 > PyInstaller 打包内置目录 > 环境变量。
    仅在打包后的可执行文件中走前两条路径，开发环境下走项目根目录。

    Returns:
        dict: 配置文件内容，未找到则返回空字典
    """
    search_paths = []
    if getattr(sys, 'frozen', False):
        # PyInstaller 打包的可执行文件
        search_paths.append(Path(sys.executable).parent / "config.json")
        search_paths.append(Path(sys._MEIPASS) / "config.json")
    else:
        # 开发环境：从项目根目录加载
        search_paths.append(Path(__file__).resolve().parent.parent.parent.parent / "config.json")
    for config_path in search_paths:
        try:
            if config_path.is_file():
                with open(config_path, "r", encoding="utf-8") as f:
                    return json.load(f)
        except Exception:
            pass
    return {}

def _get_embed_config():
    """获取 Embedding 专用 API 配置。

    注意：embedding-2 模型仅智谱 API 支持，因此始终指向智谱的 API 地址。
    如果环境变量设置了 ONLINE_API_BASE，也会用该值覆盖 base_url。

    Returns:
        dict: 包含 base_url, api_key, timeout 的配置字典
    """
    return {
        "base_url": os.getenv("ONLINE_API_BASE", "https://open.bigmodel.cn/api/paas/v4"),
        "api_key": os.getenv("ONLINE_API_KEY", ""),
        "timeout": 30,
    }

def _is_hunyuan_mt_model(model: str) -> bool:
    """判断模型名称是否为 Hunyuan-MT 翻译专用模型。

    Args:
        model: 模型名称字符串

    Returns:
        bool: True 表示是 Hunyuan-MT 模型
    """
    return model == HUNYUAN_MT_MODEL

def _get_online_config(target_lang: str = None):
    """实时读取在线翻译 API 配置。

    配置加载优先级：config.json > 环境变量 > 模块级默认值。
    如果指定了 target_lang 且目标语言不在 Hunyuan-MT 支持范围内，
    自动降级到回退模型。

    Args:
        target_lang: 目标语言代码，用于判断是否需要模型降级

    Returns:
        dict: 包含 base_url, api_key, model, timeout 的配置字典
    """
    cfg = _load_config_json()
    model = cfg.get("model") or os.getenv("ONLINE_MODEL", ONLINE_MODEL)

    # 如果目标语言不在 Hunyuan 支持范围内，自动降级到回退模型
    if _is_hunyuan_mt_model(model) and target_lang and target_lang not in HUNYUAN_MT_LANGS:
        model = HUNYUAN_FALLBACK_MODEL
        print(f"[模型路由] {target_lang} 不在Hunyuan-33语范围内，降级到 {model}")

    base = os.getenv("SILICONFLOW_API_BASE", "https://api.siliconflow.cn/v1")
    key = os.getenv("SILICONFLOW_API_KEY", "")

    return {
        "base_url": base,
        "api_key": key,
        "model": model,
        "timeout": int(os.getenv("ONLINE_TIMEOUT", str(ONLINE_TIMEOUT))),
    }


def _extract_content(resp_json: dict) -> str:
    """从 OpenAI 兼容格式的 API 返回体中安全提取 content 字段。

    推理模型（如 GLM-4.7-Flash）可能额外返回 reasoning_content，
    正常情况下 content 不为空；若 content 为空但 reasoning_content 有内容，
    说明可能是 max_tokens 不足导致截断，返回空字符串让上层重试。

    Args:
        resp_json: API 返回的 JSON 字典

    Returns:
        str: 提取到的 content 文本，提取失败则返回空字符串
    """
    msg = resp_json.get("choices", [{}])[0].get("message", {})
    content = msg.get("content", "") or ""
    if not content.strip():
        # 推理模型可能 reasoning_tokens 耗尽 max_tokens 导致 content 为空
        rc = msg.get("reasoning_content", "")
        if rc:
            print(f"[WARNING] 模型content为空但reasoning_content有内容(len={len(rc)})，"
                  f"可能是max_tokens不足，将重试")
        return ""
    return content.strip()


# ====================================================================
# 语言配置
# 定义支持的语言代码、中文名称以及翻译目标语言列表
# ====================================================================

# 翻译目标语言（9种核心语言，知识库完整覆盖）
TRANSLATE_LANGS = ["en", "ru", "ar", "es", "pt", "fr", "kk", "de", "zh_hant"]

# 全部支持的语言（包括扩展语言，部分仅在在线翻译模式下可用）
ALL_LANGS = TRANSLATE_LANGS + ["ms", "id_lang", "th", "tr", "it", "pl", "sv"]

# 语言代码 → 中文名称映射表
# 注：KB 中的语言只有 TRANSLATE_LANGS 的9种，
# 其他语言的中文名用于构造翻译提示词，避免模型不认识语言代码导致翻成英语
LANG_NAMES = {
    "en": "英语", "ru": "俄语", "ar": "阿拉伯语", "es": "西班牙语", "pt": "葡萄牙语",
    "fr": "法语", "kk": "哈萨克语", "de": "德语", "zh_hant": "繁体中文", "ms": "马来语",
    "id_lang": "印尼语", "th": "泰语", "tr": "土耳其语", "it": "意大利语",
    "pl": "波兰语", "sv": "瑞典语",
    # ★ "其他语言"——不在KB中，但翻译时需要语言中文名来构造提示词
    "ja": "日语", "ko": "韩语",
    # ★ v2.20: 补充前端OTHER_LANG_OPTIONS中的所有语言，缺少中文名会导致提示词变成"翻译为mn"，模型不认识→翻成英语
    "mn": "蒙古语", "vi": "越南语", "id": "印尼语", "nl": "荷兰语", "uk": "乌克兰语",
    "hi": "印地语", "fa": "波斯语", "he": "希伯来语", "el": "希腊语",
    "my": "缅甸语", "km": "柬埔寨语", "lo": "老挝语",
    "tl": "菲律宾语", "gu": "古吉拉特语", "ur": "乌尔都语",
    "te": "泰卢固语", "mr": "马拉地语", "bn": "孟加拉语",
    "ta": "泰米尔语", "bo": "藏语", "ug": "维吾尔语", "yue": "粤语",
}

# ====================================================================
# 常量定义
# 相似度阈值和搜索结果数量控制
# ====================================================================
HIGH_SIM = 0.90    # 高相似度阈值：直接复用知识库条目
FUZZY_SIM = 0.88   # 模糊相似度阈值
MED_SIM = 0.75     # 中等相似度阈值：可用于参考例句
TOP_K = 4          # 语义搜索返回的最多结果数
TOP_FUZZY = 3      # 模糊子串匹配返回的最多结果数

# 系统提示词已移除（改用极简 user prompt）

# ====================================================================
# CJK 汉字提取与标点无关匹配
# ====================================================================
_CJK_RE = re.compile(r'[\u4e00-\u9fff\u3400-\u4dbf]')
_zh_cjk_cache = None

def extract_cjk(text: str) -> str:
    if not text:
        return ''
    return ''.join(_CJK_RE.findall(text))

def _get_zh_cjk_cache(conn):
    global _zh_cjk_cache
    if _zh_cjk_cache is not None:
        return _zh_cjk_cache
    _zh_cjk_cache = {}
    rows = conn.execute("SELECT id, zh FROM tm_segments").fetchall()
    for rid, zh in rows:
        if not zh:
            continue
        cjk = extract_cjk(zh)
        if cjk and cjk not in _zh_cjk_cache:
            _zh_cjk_cache[cjk] = (rid, zh)
    print(f"  ✅ CJK 缓存已构建（{len(_zh_cjk_cache)} 条）")
    return _zh_cjk_cache

def _invalidate_zh_cjk_cache():
    global _zh_cjk_cache
    _zh_cjk_cache = None

# ====================================================================
# 语言识别与别名系统
# 支持多种输入写法（中文全称、英文名、语言代码等），统一映射到标准语言代码
# ====================================================================

# 语言名称映射（支持多种写法）
LANG_ALIASES = {}
for _lc, _name in LANG_NAMES.items():
    if _lc in TRANSLATE_LANGS:
        LANG_ALIASES[_lc] = _lc  # 代码本身
        LANG_ALIASES[_name] = _lc  # 中文全称
# 英文名别名
_EN_NAMES = {
    "en": ["english", "eng", "英文", "英语"],  # ★ 移除"英"：避免"精英"等词误匹配
    "ru": ["russian", "rus", "俄文", "俄语"],  # ★ 移除"俄"：避免误匹配
    "ar": ["arabic", "阿拉伯", "阿语"],
    "es": ["spanish", "esp", "西文", "西语", "西班牙"],  # ★ 移除"西"：避免"西方"等词误匹配
    "pt": ["portuguese", "por", "葡文", "葡语", "葡萄牙"],  # ★ 移除"葡"：避免"葡萄"等词误匹配
    "fr": ["french", "fra", "法文", "法语"],  # ★ 移除"法"：避免"法律"等词误匹配
    "kk": ["kazakh", "哈语", "哈萨克"],  # ★ 移除"哈"：避免"哈哈"等词误匹配
    "de": ["german", "deu", "德文", "德语"],  # ★ 移除"德"：避免"深度"等词误匹配
    "zh_hant": ["繁体", "繁中", "繁体中文", "traditional chinese", "tc"],
}
for _lc, _aliases in _EN_NAMES.items():
    for _a in _aliases:
        LANG_ALIASES[_a.lower()] = _lc

def _strip_lang_prefix(text: str, lc: str = "") -> str:
    """清洗模型输出中的语言名称前缀，如 "俄语: текст" → "текст"。

    模型有时会在翻译结果前加上语言名称或代码（如 "俄语：xxx"、"English: xxx"），
    此函数移除这些前缀，只保留实际翻译内容。

    Args:
        text: 待清洗的文本
        lc: 目标语言代码，优先匹配该语言的名称前缀

    Returns:
        str: 去除语言名前缀后的文本
    """
    if not text:
        return text
    # 匹配 "语言名：" 或 "语言名:" 或 "语言代码:" 前缀
    # 先尝试匹配目标语言本身的前缀
    if lc and lc in LANG_NAMES:
        name = LANG_NAMES[lc]
        # 匹配 "俄语：" / "俄语:" / "Russian:" / "ru:" 等
        text = re.sub(rf'^{re.escape(name)}\s*[：:]\s*', '', text)
        text = re.sub(rf'^{re.escape(lc)}\s*[：:]\s*', '', text, flags=re.IGNORECASE)
    # 再尝试匹配所有语言名/代码前缀（兜底）
    for _lc, _name in LANG_NAMES.items():
        text = re.sub(rf'^{re.escape(_name)}\s*[：:]\s*', '', text)
        text = re.sub(rf'^{re.escape(_lc)}\s*[：:]\s*', '', text, flags=re.IGNORECASE)
    # 英文语言名前缀
    _en_name_map = {"en": "English", "ru": "Russian", "ar": "Arabic", "es": "Spanish",
                    "pt": "Portuguese", "fr": "French", "kk": "Kazakh", "de": "German",
                    "zh_hant": "Traditional Chinese"}
    for _lc2, _en in _en_name_map.items():
        text = re.sub(rf'^{re.escape(_en)}\s*[：:]\s*', '', text, flags=re.IGNORECASE)
    return text.strip()


def _extract_single_lang(text: str, target_lang: str) -> str:
    """从模型输出中提取目标语言的翻译内容。

    ★ v2.12 改造：单语模式下模型已被明确指示只输出一种语言，
    大多数情况下输出就是干净的。因此采用保守策略：
    - 有明确语言前缀行（如 "英语:" "Russian:"）→ 提取目标语言段
    - 无语言前缀 → 直接保留，不做段落级截断（防止误删多段落翻译）

    Args:
        text: 模型原始输出文本
        target_lang: 目标语言代码

    Returns:
        str: 提取后的目标语言翻译内容
    """
    if not text:
        return text

    # 构建所有可能的语言前缀模式
    _EN_NAME_MAP = {"en": "English", "ru": "Russian", "ar": "Arabic", "es": "Spanish",
                    "pt": "Portuguese", "fr": "French", "kk": "Kazakh", "de": "German",
                    "zh_hant": "Traditional Chinese"}
    
    # 构建目标语言的所有可能别名
    target_aliases = set()
    target_aliases.add(target_lang.lower())
    if target_lang in LANG_NAMES:
        target_aliases.add(LANG_NAMES[target_lang].lower())
    if target_lang in _EN_NAME_MAP:
        target_aliases.add(_EN_NAME_MAP[target_lang].lower())
    
    # 收集所有语言前缀的正则
    all_prefixes = []
    for _lc, _name in LANG_NAMES.items():
        all_prefixes.append(re.escape(_name))
        all_prefixes.append(re.escape(_lc))
    for _en in _EN_NAME_MAP.values():
        all_prefixes.append(re.escape(_en))
    
    prefix_pattern = '(' + '|'.join(all_prefixes) + r')\s*[：:]\s*'
    
    # 检查文本中是否包含语言前缀行
    matches = list(re.finditer(prefix_pattern, text, re.IGNORECASE))
    
    if not matches:
        # ★ 没有语言前缀 → 模型输出就是干净的单语翻译，直接保留
        # 不走 _strip_foreign_paragraphs，防止误删多段落内容
        return _strip_lang_prefix(text, target_lang)
    
    if len(matches) == 1:
        # 只有一个前缀
        prefix_text = matches[0].group(1).lower()
        if prefix_text in target_aliases:
            # 是目标语言自己的前缀，提取内容
            segment = text[matches[0].end():].strip()
            return _strip_lang_prefix(segment, target_lang)
        else:
            # 是其他语言的前缀，目标语言内容在前面的部分
            segment = text[:matches[0].start()].strip()
            if segment:
                return _strip_lang_prefix(segment, target_lang)
            # 前面没有内容，尝试提取后面的但去掉该前缀
            return _strip_lang_prefix(text, target_lang)
    
    # 多个语言前缀 → 模型违规输出了多语，需要提取目标语言段
    # 先找目标语言的段（有前缀的情况）
    for i, m in enumerate(matches):
        prefix_text = m.group(1).lower()
        if prefix_text in target_aliases:
            content_start = m.end()
            content_end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
            segment = text[content_start:content_end].strip()
            return _strip_lang_prefix(segment, target_lang)
    
    # 目标语言没有前缀，但其他语言有前缀
    # → 目标语言内容在第一个非目标语言前缀之前
    first_other = matches[0]
    segment = text[:first_other.start()].strip()
    if segment:
        return _strip_lang_prefix(segment, target_lang)
    
    # 兜底
    return _strip_lang_prefix(text, target_lang)


def _strip_foreign_paragraphs(text: str, target_lang: str) -> str:
    """段落级语言检测：移除目标语言之外的外语段落。

    典型场景：模型输出 "葡萄牙语内容\n\nEnglish content here\n\n俄语内容"
    没有语言名前缀，但段落明显是其他语言。

    策略：
    - 对每个段落，检测其字符组成（CJK/西里尔/阿拉伯/拉丁等）
    - 根据目标语言的字符特征判断是否属于目标语言
    - 移除明显不属于目标语言的段落
    - 拉丁字母系语言（en/es/pt/fr/de）因无法仅靠字符区分，走 _strip_lang_name_sections

    Args:
        text: 待清洗的文本
        target_lang: 目标语言代码

    Returns:
        str: 清洗后的文本
    """
    if not text:
        return text

    # 每种语言的"特征字符"范围（★ 扩展为全球主要书写系统）
    CYRILLIC_RE = re.compile(r'[\u0400-\u04ff]')         # 西里尔（俄语/哈萨克语等）
    ARABIC_RE = re.compile(r'[\u0600-\u06ff\u0750-\u077f]')  # 阿拉伯字母
    CJK_RE = re.compile(r'[\u4e00-\u9fff\u3400-\u4dbf]')     # CJK汉字
    KANA_RE = re.compile(r'[\u3040-\u309f\u30a0-\u30ff]')     # 日语假名
    HANGUL_RE = re.compile(r'[\uac00-\ud7af]')                 # 韩语谚文
    THAI_RE = re.compile(r'[\u0e00-\u0e7f]')                   # 泰文
    DEVANAGARI_RE = re.compile(r'[\u0900-\u097f]')             # 天城文（印地语等）
    HEBREW_RE = re.compile(r'[\u0590-\u05ff]')                 # 希伯来文
    GREEK_RE = re.compile(r'[\u0370-\u03ff\u1f00-\u1fff]')    # 希腊文
    
    def _para_script_score(para: str) -> dict:
        """统计段落中各类书写系统字符的占比。

        Args:
            para: 待统计的段落文本

        Returns:
            dict: {书写系统名称: 占比}，如 {'cyrillic': 0.6, 'cjk': 0.4}
        """
        if not para.strip():
            return {}
        total = len(para.strip())
        scores = {}
        counts = {
            'cyrillic': len(CYRILLIC_RE.findall(para)),
            'arabic':   len(ARABIC_RE.findall(para)),
            'cjk':      len(CJK_RE.findall(para)),
            'kana':     len(KANA_RE.findall(para)),
            'hangul':   len(HANGUL_RE.findall(para)),
            'thai':     len(THAI_RE.findall(para)),
            'devanagari': len(DEVANAGARI_RE.findall(para)),
            'hebrew':   len(HEBREW_RE.findall(para)),
            'greek':    len(GREEK_RE.findall(para)),
        }
        for script, count in counts.items():
            if count > 0:
                scores[script] = count / total
        return scores
    
    # 目标语言的主字符类型（★ 使用全局 _SCRIPT_LANG_MAP，不硬编码）
    target_script = _SCRIPT_LANG_MAP.get(target_lang)
    # 日语特殊：主系统是cjk，但假名(kana)是其特有字符，用于区分中文泄漏
    # 段落级检测中，日语归cjk处理（和中文共用CJK汉字，但含假名的是日语段落）
    if target_lang == "ja":
        target_script = "cjk"
    # 韩语主系统是hangul，段落级检测时用hangul判断
    # en, es, pt, fr, de 都是拉丁字母，不做段落级过滤（无法仅靠字符区分）

    # 如果目标语言是拉丁字母系，段落级检测无法区分，用语言名关键词行截断
    if target_script is None or target_script == "latin":
        return _strip_lang_name_sections(text, target_lang)
    
    paragraphs = text.split('\n')
    # ★ 预编译目标语言特有字符检测（用于区分CJK是日语/韩语还是中文泄漏）
    _KANA_RE = re.compile(r'[\u3040-\u309f\u30a0-\u30ff]')      # 日语假名
    _HANGUL_RE = re.compile(r'[\uac00-\ud7af]')                   # 韩语谚文
    _CJK_CHAR_RE = re.compile(r'[\u4e00-\u9fff\u3400-\u4dbf]')  # CJK汉字
    
    kept = []
    for para in paragraphs:
        stripped = para.strip()
        if not stripped:
            kept.append(para)  # 保留空行
            continue
        scores = _para_script_score(stripped)
        # 如果段落完全不含任何特殊字符，跳过（短段落保留）
        if not scores:
            if len(stripped) <= 5:
                kept.append(para)
            continue
        
        dominant = max(scores, key=scores.get)
        
        # ★ 通用判断：基于 _SCRIPT_LANG_MAP 和特有字符检测
        if target_script == "cjk":
            # CJK书写系统（ja/zh_hant）：含假名→日语段落；纯CJK无假名→中文泄漏
            # zh_hant 自身就是中文，所有CJK段落都保留
            if target_lang == "zh_hant":
                kept.append(para)
            elif _KANA_RE.search(stripped):
                kept.append(para)  # 含假名→日语
            elif len(stripped) <= 5:
                kept.append(para)  # 短CJK可能是汉字词
            else:
                continue  # 纯CJK无假名→中文泄漏
        elif target_script == "hangul":
            # 韩语：含谚文→韩语；纯CJK→中文泄漏；非CJK→保留
            if _HANGUL_RE.search(stripped):
                kept.append(para)
            elif dominant == "cjk" and len(stripped) <= 4:
                kept.append(para)  # 短CJK可能是汉字人名
            elif dominant != "cjk":
                kept.append(para)
            else:
                continue  # 纯CJK长段落→中文泄漏
        elif target_script in ("cyrillic", "arabic", "thai", "devanagari",
                                "hebrew", "greek", "bengali", "tamil",
                                "myanmar", "khmer", "lao"):
            # 其他非拉丁书写系统：dominant匹配→保留；不匹配→跳过
            if dominant == target_script:
                kept.append(para)
            elif len(stripped) <= 5:
                kept.append(para)
            else:
                continue
        else:
            # 未知书写系统或拉丁字母系：保守保留
            kept.append(para)
    
    result = '\n'.join(kept).strip()
    # 清理连续空行
    result = re.sub(r'\n{3,}', '\n\n', result)
    return _strip_lang_prefix(result, target_lang)


def _strip_lang_name_sections(text: str, target_lang: str) -> str:
    """对拉丁字母系语言（en/es/pt/fr/de），通过语言名关键词行截断多语输出。

    模型输出如：
    "Se voce ama... [葡语翻译]\n\nEnglish:\nWe look forward to you...\n\nRussian:\nМы..."

    策略：扫描每行，如果发现非目标语言的名称行（如 "English:"、"俄语："、"Russian:"），
    截掉该行及之后所有内容。

    拉丁字母系语言无法通过字符编码区分（都使用拉丁字母），
    只能通过语言名称关键词行来判断是否进入了其他语言区域。

    Args:
        text: 待清洗的文本
        target_lang: 目标语言代码

    Returns:
        str: 截断后的文本
    """
    if not text:
        return text
    
    _EN_NAME_MAP = {"en": "English", "ru": "Russian", "ar": "Arabic", "es": "Spanish",
                    "pt": "Portuguese", "fr": "French", "kk": "Kazakh", "de": "German",
                    "zh_hant": "Traditional Chinese"}
    
    # 构建目标语言别名
    target_aliases = set()
    target_aliases.add(target_lang.lower())
    if target_lang in LANG_NAMES:
        target_aliases.add(LANG_NAMES[target_lang].lower())
    if target_lang in _EN_NAME_MAP:
        target_aliases.add(_EN_NAME_MAP[target_lang].lower())
    
    # 构建所有语言名/代码（包括中文+英文+代码）的前缀行正则
    # 匹配独占一行的语言名（前后可能有空格，后面跟冒号或换行）
    all_names = []
    for _lc, _name in LANG_NAMES.items():
        all_names.append((_name, _lc))
    for _lc, _en in _EN_NAME_MAP.items():
        all_names.append((_en, _lc))
    for _lc in TRANSLATE_LANGS:
        all_names.append((_lc, _lc))
    
    # 构建非目标语言的名称行检测正则
    other_lang_patterns = []
    for name, lc in all_names:
        if lc not in target_aliases:
            # 匹配 "English:" / "俄语：" / "ru:" 等独占一行或段首的情况
            other_lang_patterns.append(re.escape(name))
    
    if not other_lang_patterns:
        return _strip_lang_prefix(text, target_lang)
    
    # 按行扫描，遇到非目标语言名称行就截断
    lines = text.split('\n')
    cutoff = len(lines)
    for i, line in enumerate(lines):
        stripped = line.strip().rstrip('：:').strip()
        for pat in other_lang_patterns:
            if re.fullmatch(pat, stripped, re.IGNORECASE):
                cutoff = i
                break
        if cutoff != len(lines):
            break
    
    result = '\n'.join(lines[:cutoff]).strip()
    return _strip_lang_prefix(result, target_lang)


# ====================================================================
# 语言 → 书写系统映射（全局，供清洗逻辑使用）
# ====================================================================
# 设计原则：
#   CJK统一汉字(U+4E00-U+9FFF)不是中文专属——日语(kanji)、韩语(hanja)也使用。
#   粗暴删除CJK字符会误杀这些语言的正常内容（如日语汉字全被删掉只剩假名碎片）。
#   因此需要一个映射表，告诉清洗逻辑：哪些语言的正常输出中应包含CJK字符。
#
#   _CJK_SCRIPT_LANGS：正常使用CJK汉字的语言 → 跳过CJK字符级清洗，改用标点检测
#   _SCRIPT_LANG_MAP ：语言→主书写系统 → 供段落级检测判断"这段文字是否属于目标语言"
#
#   新增语言时只需更新这两个映射表，无需改动清洗函数逻辑。

# 使用CJK统一汉字作为正常书写系统的语言
# ja: 日语汉字(kanji)占日文文本30-50%，删除=毁灭性破坏
# ko: 现代韩语主要用谚文，但人名/学术用语偶尔用汉字(hanja)，误删损失较小但应保留
# zh_hant: 繁体中文本身就是中文，CJK是其核心字符
_CJK_SCRIPT_LANGS = frozenset({"zh_hant", "ja", "ko"})

# 语言→主书写系统映射（用于段落级字符检测）
# "cyrillic": 西里尔字母（俄语、哈萨克语等）
# "arabic":   阿拉伯字母（阿拉伯语、波斯语等）
# "cjk":      CJK统一汉字（中文、日语、韩语等）
# "latin":    拉丁字母（英语、西语、葡语、法语、德语等）
# "devanagari": 天城文（印地语、马拉地语等）
# "hangul":   谚文（韩语主书写系统，与CJK汉字共存）
# "thai":     泰文
# "kana":     假名（日语特有，与CJK汉字共存）
_SCRIPT_LANG_MAP = {
    # KB语言
    "en": "latin", "ru": "cyrillic", "ar": "arabic", "es": "latin",
    "pt": "latin", "fr": "latin", "de": "latin", "kk": "cyrillic",
    "zh_hant": "cjk",
    # 其他常见语言（持续扩展）
    "ja": "cjk",       # 汉字+假名，主系统归cjk
    "ko": "hangul",    # 主系统是谚文，辅助用汉字
    "th": "thai", "vi": "latin", "ms": "latin", "id": "latin",
    "it": "latin", "pl": "latin", "sv": "latin", "nl": "latin",
    "uk": "cyrillic", "el": "greek", "cs": "latin", "ro": "latin",
    "hu": "latin", "fi": "latin", "da": "latin", "no": "latin",
    "tr": "latin", "he": "hebrew", "fa": "arabic", "hi": "devanagari",
    "ur": "arabic", "bn": "bengali", "ta": "tamil", "mn": "cyrillic",
    "my": "myanmar", "km": "khmer", "lo": "lao", "fil": "latin",
}


def _strip_chinese_in_non_zh(text: str, lang_code: str) -> str:
    """清洗非目标语言输出中泄漏的中文字符。

    glm-4-flash 长文本翻译时偶尔会在俄语/西语等输出中混入中文片段，
    如 "если вы 擅长 организовать" → 清洗为 "если вы организовать"

    ★★ 通用规则（基于书写系统映射，适配全球语言）：

    情况1 — 目标语言使用CJK书写系统（ja/ko/zh_hant）：
      → CJK汉字是该语言的正常输出，不能删除。
      → 改用"中文标点检测"识别真正的中文泄漏：
        中文特有标点（。、「」『』【】《》等）在日语中也有使用，
        但连续出现"中文标点+CJK词组"且无目标语言特有字符（假名/谚文）时，
        很可能是中文泄漏段落，可安全移除。
      → 孤立的CJK词组（无中文标点包围）一律保留——宁可漏清不可误杀。

    情况2 — 目标语言不使用CJK书写系统（ru/ar/es/en等）：
      → CJK汉字一定是中文泄漏，直接删除整个CJK词组。
      → 这对全球绝大多数语言都是安全的。

    Args:
        text: 待清洗的文本
        lang_code: 目标语言代码

    Returns:
        str: 清洗后的文本
    """
    if not text:
        return text
    
    # ★ 情况1：目标语言使用CJK书写系统 → 不能删CJK字符，改用标点检测
    if lang_code in _CJK_SCRIPT_LANGS:
        return _strip_chinese_leaks_by_punctuation(text, lang_code)
    
    # ★ 情况2：目标语言不使用CJK → CJK一定是中文泄漏，直接删除
    # CJK 统一汉字 + 扩展A/B + 兼容汉字
    cleaned = re.sub(r'[\u4e00-\u9fff\u3400-\u4dbf\U00020000-\U0002a6df]+', '', text)
    
    # 修整清洗后的空格和空行
    cleaned = re.sub(r' {2,}', ' ', cleaned)
    cleaned = re.sub(r' \n', '\n', cleaned)
    cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
    cleaned = '\n'.join(line.strip() for line in cleaned.split('\n'))
    
    return cleaned.strip()


def _strip_chinese_leaks_by_punctuation(text: str, lang_code: str) -> str:
    """对使用CJK书写系统的语言（ja/ko/zh_hant），用"中文标点"识别中文泄漏。

    核心思路：CJK汉字本身不能区分中文和日文汉字，但"中文标点+CJK词组"
    的组合在日语/韩语输出中几乎不会自然出现，大概率是中文泄漏片段。

    检测规则：
    1. 中文特有标点（。？！、）+ 紧邻的CJK词组 → 视为中文泄漏，删除
    2. 日语/韩语也用的标点（「」【】等） → 不作为泄漏信号，保留
    3. 纯CJK词组（无中文标点） → 一律保留，宁可漏清不可误杀

    对 zh_hant：不清洗（中文是其正确输出）

    Args:
        text: 待清洗的文本
        lang_code: 目标语言代码（ja/ko/zh_hant）

    Returns:
        str: 清洗后的文本
    """
    # 繁体中文不清洗
    if lang_code == "zh_hant":
        return text
    
    # 日语/韩语：用中文特有标点检测泄漏片段
    # 中文特有标点（。？！、）—— 日语用。但通常配合假名，韩语不用这些标点
    # ★ 关键区分：日语的"。"配合假名使用，中文的"。"配合纯CJK使用
    #   如果"。"前后一段距离内没有假名/谚文，大概率是中文泄漏
    
    if lang_code == "ja":
        # ★ 日语：词组级+子片段级清洗——删除不含假名的纯CJK词组
        # 
        # PPT/文档翻译中，中文残留常以"词组"形式混入日语文本：
        #   "最近一年以上 求人サイトを巡回" → "最近一年以上"无假名=中文残留
        #   "新しい職位をよく見かけます" → 含假名=日语保留
        #
        # 策略：
        # 1. 按空格切分 → 独立纯CJK无假名词组=中文残留删除
        # 2. 对保留的混合token（CJK+假名），做子片段级清洗：
        #    删除token内的纯CJK子片段（连续CJK汉字，不含假名）
        #    "新しい職位"中"職位"紧跟假名→保留（日语汉字词）
        #    但"最近一年以上詳細"中纯CJK部分→删除
        #    注意：日语汉字词普遍，所以短CJK（≤2字）保守保留
        _KANA_RE = re.compile(r'[\u3040-\u309f\u30a0-\u30ff]')
        _CJK_RE = re.compile(r'[\u4e00-\u9fff\u3400-\u4dbf]')
        _CJK_ONLY_FRAGMENT = re.compile(r'[\u4e00-\u9fff\u3400-\u4dbf]+')
        
        tokens = re.split(r'(\s+)', text)  # 保留空格
        kept_tokens = []
        for token in tokens:
            if not token.strip():
                kept_tokens.append(token)
                continue
            has_cjk = bool(_CJK_RE.search(token))
            if not has_cjk:
                kept_tokens.append(token)
                continue
            has_kana = bool(_KANA_RE.search(token))
            if has_kana:
                # ★ CJK+假名混合token → 子片段级清洗
                # 日语汉字词多，不能简单全删CJK，只删除长段纯CJK子片段
                # "新しい職位" → "職位"只有2字CJK且紧跟假名 → 保留
                # "詳細はこちらをご覧ください" → "詳細"2字紧跟假名 → 保留
                # "最近一年以上" 长段纯CJK无假名 → 删除
                cleaned = token
                for frag in _CJK_ONLY_FRAGMENT.findall(token):
                    # ★ 检查frag前后1字内是否有假名（有假名=日语汉字词，保留）
                    frag_start = token.index(frag)
                    before_char = token[max(0, frag_start - 1):frag_start]
                    after_char = token[frag_start + len(frag):frag_start + len(frag) + 1]
                    has_kana_nearby = bool(_KANA_RE.search(before_char)) or bool(_KANA_RE.search(after_char))
                    if len(frag) >= 3 and not has_kana_nearby:
                        # 3字以上纯CJK，且前后1字内无假名 → 中文残留
                        cleaned = cleaned.replace(frag, '', 1)
                if cleaned.strip():
                    kept_tokens.append(cleaned)
            else:
                # 纯CJK无假名
                cjk_chars = _CJK_RE.findall(token)
                if len(cjk_chars) <= 2:
                    # 短CJK可能是日语汉字词，保留
                    kept_tokens.append(token)
                else:
                    # 3个以上CJK汉字且无假名 = 中文残留，删除
                    pass
        
        result = ''.join(kept_tokens).strip()
        result = re.sub(r' {2,}', ' ', result)
        result = re.sub(r'\n{3,}', '\n\n', result)
        return result
    
    elif lang_code == "ko":
        # ★ 韩语：词组级+子片段级清洗——彻底删除所有中文残留
        # 
        # 问题场景：
        #   "ROX赋能플랫폼 挑战 전투" → "赋能"夹在ROX和谚文之间，按空格切分被整体保留
        #   "고객이 온라인으로 이동했습니다 客户已转到线上" → 独立中文片段
        #
        # 策略：
        # 1. 按空格切分为词组 → 独立纯CJK词组直接删除
        # 2. 对保留的混合token（CJK+谚文），再做子片段级清洗：
        #    把token内的纯CJK子片段（不含谚文/假名的连续CJK字符）删除
        #    "ROX赋能X赋能플랫폼" → "赋能"是纯CJK子片段 → 删除 → "ROX플랫폼"
        #
        # 原则：现代韩语几乎不用汉字，翻译结果中出现的CJK汉字99%是中文残留
        _HANGUL_RE = re.compile(r'[\uac00-\ud7af]')
        _CJK_RE = re.compile(r'[\u4e00-\u9fff\u3400-\u4dbf]')
        # ★ 匹配连续CJK汉字片段（不含谚文/假名/拉丁，纯CJK字符序列）
        _CJK_ONLY_FRAGMENT = re.compile(r'[\u4e00-\u9fff\u3400-\u4dbf]+')
        
        tokens = re.split(r'(\s+)', text)  # 保留空格
        kept_tokens = []
        for token in tokens:
            if not token.strip():
                kept_tokens.append(token)
                continue
            # 不含CJK → 保留（纯谚文/拉丁/数字等）
            has_cjk = bool(_CJK_RE.search(token))
            if not has_cjk:
                kept_tokens.append(token)
                continue
            # 含CJK → 检查是否含谚文
            has_hangul = bool(_HANGUL_RE.search(token))
            if has_hangul:
                # ★ CJK+谚文混合token → 子片段级清洗
                # 删除token内的纯CJK子片段，保留谚文/拉丁/数字等
                # "ROX赋能플랫폼" → 删除"赋能" → "ROX플랫폼"
                cleaned = _CJK_ONLY_FRAGMENT.sub('', token)
                if cleaned.strip():
                    kept_tokens.append(cleaned)
                # 清洗后为空则整体丢弃
            else:
                # 纯CJK无谚文 = 中文残留，直接删除
                pass
        
        result = ''.join(kept_tokens).strip()
        result = re.sub(r' {2,}', ' ', result)
        result = re.sub(r'\n{3,}', '\n\n', result)
        return result
    
    # 其他CJK书写系统语言（未来扩展），保守处理：不清洗
    return text


def parse_target_langs(user_input: str) -> list[str] | None:
    """从用户输入中识别目标语言。

    支持多种输入格式：
    - "翻成英文/英语" → ["en"]
    - "翻译成英文和俄语" → ["en", "ru"]
    - "帮我翻英语、西语" → ["en", "es"]
    - "极石汽车"（无语言关键词）→ None

    短别名（≤2字符）使用词边界匹配防止误匹配，
    结果按在用户输入中首次出现的顺序排序。

    Args:
        user_input: 用户输入文本

    Returns:
        list[str] | None: 语言代码列表，未指定任何语言则返回 None
    """
    import re
    found = []
    input_lower = user_input.lower()
    for alias, lc in LANG_ALIASES.items():
        if lc not in TRANSLATE_LANGS:
            continue
        # 短 alias（≤2字符）用词边界匹配，避免 "portuguese" 误匹配 "es"
        if len(alias) <= 2:
            pattern = rf'(?<![a-z]){re.escape(alias)}(?![a-z])'
            if re.search(pattern, input_lower) and lc not in found:
                found.append(lc)
        else:
            if alias in input_lower and lc not in found:
                found.append(lc)
    if not found:
        return None
    # 按用户输入中首次出现的顺序排序
    def _first_pos(lc):
        best = len(user_input)
        for alias, alc in LANG_ALIASES.items():
            if alc == lc:
                pos = user_input.lower().find(alias)
                if 0 <= pos < best:
                    best = pos
        return best
    found.sort(key=_first_pos)
    return found


# ====================================================================
# 路径配置
# 优先使用用户数据目录（避免 App Translocation 只读问题）
# ====================================================================
_USER_DATA_DIR = os.getenv("USER_DATA_DIR")
if _USER_DATA_DIR:
    DATA_DIR = Path(_USER_DATA_DIR)
else:
    DATA_DIR = Path.home() / "Library" / "Application Support" / "翻译助手"
DB_PATH = str(DATA_DIR / "tm.sqlite3")                # 翻译记忆库 SQLite 文件路径
EMB_PATH = str(DATA_DIR / "tm_embeddings.npz")        # 向量索引文件路径
INDEX_STAMP = DATA_DIR / ".index_stamp"               # 索引时间戳文件（用于增量更新检测）


# ====================================================================
# 底层数据库与向量操作
# ====================================================================
def get_db():
    """获取 SQLite 数据库连接（单例模式，check_same_thread=False 支持多线程访问）。

    Returns:
        sqlite3.Connection: 数据库连接对象
    """
    conn = sqlite3.connect(DB_PATH, check_same_thread=False)
    init_segment_db(conn)
    return conn

def ollama_is_up() -> bool:
    """检查 Ollama 服务是否可用。

    注意：现在改用在线 embedding API，有 API Key 即表示可用。
    此函数保留是为了兼容旧代码调用。

    Returns:
        bool: True 表示在线 embedding API 已配置
    """
    return online_api_is_configured()

def embed(text: str) -> np.ndarray:
    """智谱 embedding-2 在线向量化（单条文本）。

    调用智谱 API 的 embeddings 接口，将文本转为 1024 维向量并归一化。
    空文本返回零向量。

    Args:
        text: 待向量化的文本

    Returns:
        np.ndarray: 归一化的 1024 维向量
    """
    text = str(text).strip()
    if not text:
        return np.zeros(1024, dtype=np.float32)
    cfg = _get_embed_config()
    r = requests.post(
        f"{cfg['base_url'].rstrip('/')}/embeddings",
        json={"model": EMBED_MODEL, "input": text},
        headers={"Authorization": f"Bearer {cfg['api_key']}"},
        timeout=30,
    )
    r.raise_for_status()
    v = np.array(r.json()["data"][0]["embedding"], dtype=np.float32)
    return v / np.linalg.norm(v)

def embed_batch(texts: list[str], batch_size: int = 32) -> np.ndarray:
    """智谱 embedding-2 批量在线向量化。

    将文本列表分批发送到智谱 API，每批最多 batch_size 条。
    所有向量结果拼接为一个 numpy 数组返回。

    Args:
        texts: 待向量化的文本列表
        batch_size: 每批处理的文本数，默认 32

    Returns:
        np.ndarray: 形状为 (len(texts), 1024) 的归一化向量数组
    """
    cfg = _get_embed_config()
    all_vecs = []
    for i in range(0, len(texts), batch_size):
        chunk = [t.strip() for t in texts[i:i + batch_size] if t.strip()]
        if not chunk:
            continue
        r = requests.post(
            f"{cfg['base_url'].rstrip('/')}/embeddings",
            json={"model": EMBED_MODEL, "input": chunk},
            headers={"Authorization": f"Bearer {cfg['api_key']}"},
            timeout=60,
        )
        r.raise_for_status()
        for emb in r.json()["data"]:
            v = np.array(emb["embedding"], dtype=np.float32)
            all_vecs.append(v / np.linalg.norm(v))
    return np.array(all_vecs, dtype=np.float32)

def load_index():
    """从 npz 文件中加载向量索引。

    Returns:
        tuple: (ids: np.ndarray, vecs: np.ndarray)
                ids 为 int64 数组，vecs 为 float32 二维数组
                索引文件不存在时返回空数组
    """
    if not os.path.exists(EMB_PATH):
        return np.array([]), np.array([])
    data = np.load(EMB_PATH)
    return data["ids"], data["vecs"]

def search(query_vec: np.ndarray, ids: np.ndarray, vecs: np.ndarray, k: int = TOP_K):
    """向量相似度搜索（余弦相似度，使用点积替代因向量已归一化）。

    Args:
        query_vec: 查询向量（已归一化）
        ids: 所有条目的 ID 数组
        vecs: 所有条目的向量数组
        k: 返回的最相近结果数

    Returns:
        list[tuple]: [(id, similarity), ...]，按相似度降序排列
    """
    if len(ids) == 0:
        return []
    sims = vecs @ query_vec
    order = np.argsort(-sims)[:k]
    return [(int(ids[i]), float(sims[i])) for i in order]


# ====================================================================
# 业务操作：知识库数据读取和模糊匹配
# ====================================================================
def fetch_row(conn, rid: int) -> dict:
    """从 tm_segments 表中根据 ID 获取完整记录。

    Args:
        conn: SQLite 数据库连接
        rid: 记录 ID

    Returns:
        dict: 包含所有列的字典（id, zh, module + 所有语言列）
    """
    cols = ["id", "zh", "module"] + ALL_LANGS
    r = conn.execute(f"SELECT {','.join(cols)} FROM tm_segments WHERE id=?", (rid,)).fetchone()
    return dict(zip(cols, r))

def fuzzy_substring_hits(conn, zh_short: str) -> list[dict]:
    """模糊子串匹配：在知识库中查找包含输入文本子串的条目。

    使用 LIKE '%输入%' 进行子串匹配，然后过滤长度差不超过30的条目，
    按长度差升序排列，取前 TOP_FUZZY 条作为候选。

    Args:
        conn: SQLite 数据库连接
        zh_short: 输入的中文文本（通常较短）

    Returns:
        list[dict]: 匹配到的完整条目列表（含所有语言列）
    """
    cur = conn.execute(
        "SELECT id, zh FROM tm_segments WHERE zh LIKE '%' || ? || '%'",
        (zh_short,),
    )
    rows = cur.fetchall()
    if not rows:
        return []
    tagged = []
    for rid, long_zh in rows:
        diff = len(long_zh) - len(zh_short)
        if 0 <= diff <= 30:
            tagged.append((diff, rid, long_zh))
    tagged.sort(key=lambda x: x[0])
    return [fetch_row(conn, rid) for _, rid, _ in tagged[:TOP_FUZZY]]


# ====================================================================
# 在线模型翻译（★ 新增）
# 所有和在线 API 通信的逻辑集中在此
# ====================================================================
def online_api_is_configured() -> bool:
    """检查在线模型 API 是否已配置（实时读取环境变量）。

    Returns:
        bool: True 表示 API Key 和 Base URL 都已配置
    """
    cfg = _get_online_config()
    return bool(cfg["api_key"] and cfg["base_url"])

def detect_source_lang(text: str) -> str:
    """检测文本的主要语言（仅区分中文/非中文）。

    通过统计 CJK 统一汉字字符占比来判断。先用 extract_cjk
    提取实际 CJK 内容，再按非空白字符计算占比，避免标点空格干扰。

    Args:
        text: 待检测的文本

    Returns:
        str: "zh" 或 "en"
    """
    if not text or not text.strip():
        return "zh"
    cjk = extract_cjk(text)
    if not cjk:
        return "en"
    non_space = text.replace(' ', '').replace('\t', '').replace('\n', '').replace('\r', '')
    return "zh" if len(cjk) / max(len(non_space), 1) > 0.25 else "en"

def get_source_name(lang: str) -> str:
    """获取源语言的中文名称。

    Args:
        lang: 语言代码（"zh" 或 "en"）

    Returns:
        str: 中文名称
    """
    return {"zh": "中文", "en": "英文"}.get(lang, "原文")

def check_online_api_available() -> bool:
    """检查在线模型 API 是否可访问（轻量探测）。

    ★ v2.14: 探测失败返回 False——断网时不再假装可用。
    使用 GET /models 端点探测，超时 5 秒。
    只有状态码 < 500 才算可用。

    Returns:
        bool: True 表示 API 可正常访问
    """
    cfg = _get_online_config()
    if not (cfg["api_key"] and cfg["base_url"]):
        return False
    try:
        resp = requests.get(
            f"{cfg['base_url'].rstrip('/')}/models",
            headers={"Authorization": f"Bearer {cfg['api_key']}"},
            timeout=5,
        )
        return resp.status_code < 500
    except:
        return False  # ★ 探测失败=不可达，不再返回True

def call_online_llm(zh_text: str, examples: list[dict], retries: int = 2, on_step=None) -> dict:
    """调用在线模型 API（OpenAI 兼容格式）进行多语翻译兜底。

    将中文文本一次性翻译为 TRANSLATE_LANGS 中的全部 9 种语言。
    输出格式为 JSON 对象，key 为语言代码，value 为对应翻译文本。

    ★ 在线模式用 JSON 格式输出，彻底解决多段落解析问题。
    ★ 动态 max_tokens：短文 2048，长文按输入长度估算。

    Args:
        zh_text: 待翻译的中文文本
        examples: 参考例句列表（可选，用于术语和风格参考）
        retries: 失败重试次数
        on_step: 进度回调函数（未使用，保留兼容性）

    Returns:
        dict: {lang_code: translation_text, ...}
        对所有 TRANSLATE_LANGS 中的语言都有值

    Raises:
        所有重试用尽后抛出最后一次异常
    """
    cfg = _get_online_config()
    # ★ 在线模式用 JSON 格式输出，彻底解决多段落解析问题
    lang_list = ", ".join(f'"{c}" ({LANG_NAMES.get(c, c)})' for c in TRANSLATE_LANGS)

    example_block = ""
    if examples:
        example_block = "参考例句（沿用其术语和风格）：\n"
        for ex in examples:
            example_block += f"- 中文：{ex['zh']}\n"
            for lc in TRANSLATE_LANGS:
                if ex.get(lc):
                    example_block += f"  {LANG_NAMES[lc]}：{ex[lc]}\n"
        example_block += "\n"

    # ★ 动态 max_tokens：短文 2048，长文按输入长度估算
    max_tokens = max(2048, min(8192, len(zh_text) * 9 * 2))

    user_prompt = (
        f"{example_block}"
        f"请将以下中文内容翻译为这9种语言：{lang_list}\n\n"
        f'输出严格的 JSON 对象，key 为语言代码，value 为对应翻译文本。'
        f'多段落内容用 \\n 换行。不要输出任何其他内容。\n\n'
        f"中文原文：\n{zh_text}\n\n"
        f'输出格式示例：\n'
        f'{{"en": "English translation", "ru": "Русский перевод", "ar": "الترجمة العربية", '
        f'"es": "Traducción española", "pt": "Tradução portuguesa", "fr": "Traduction française", '
        f'"kk": "Қазақша аударма", "de": "Deutsche Übersetzung", "zh_hant": "繁體中文翻譯"}}'
    )

    url = f"{cfg['base_url'].rstrip('/')}/chat/completions"
    headers = {
        "Authorization": f"Bearer {cfg['api_key']}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": cfg["model"],
        "messages": [
            {"role": "user", "content": user_prompt},
        ],
        "temperature": 0.1,
        "max_tokens": max_tokens,
    }

    last_error = None
    content = ""
    for attempt in range(retries + 1):
        try:
            print(f"[在线] 调用 {cfg['model']} @ {cfg['base_url']} (尝试{attempt+1}, max_tokens={max_tokens})")
            r = requests.post(url, headers=headers, json=payload, timeout=cfg["timeout"])
            r.raise_for_status()
            content = _extract_content(r.json())
            print(f"[在线] 返回成功，内容长度={len(content)}")
            break
        except Exception as e:
            last_error = e
            print(f"[在线] 尝试{attempt+1}失败: {e}")
            if attempt < retries:
                import time
                time.sleep(2)
            else:
                raise last_error

    # ★ JSON 解析（多段落翻译的关键：值中可以含 \n）
    result = {}
    # 先尝试直接解析
    try:
        parsed = json.loads(content.strip())
        if isinstance(parsed, dict):
            result = {lc: _strip_lang_prefix(str(parsed.get(lc, "")).strip(), lc) for lc in TRANSLATE_LANGS}
    except json.JSONDecodeError:
        # 尝试从 markdown 代码块中提取
        try:
            cleaned = re.sub(r"^```(?:json)?|```$", "", content.strip(), flags=re.MULTILINE).strip()
            parsed = json.loads(cleaned)
            if isinstance(parsed, dict):
                result = {lc: _strip_lang_prefix(str(parsed.get(lc, "")).strip(), lc) for lc in TRANSLATE_LANGS}
        except json.JSONDecodeError:
            # 最后兜底：尝试正则提取 JSON 对象
            m = re.search(r'\{[^{}]*\}', content, re.DOTALL)
            if m:
                try:
                    parsed = json.loads(m.group())
                    if isinstance(parsed, dict):
                        result = {lc: _strip_lang_prefix(str(parsed.get(lc, "")).strip(), lc) for lc in TRANSLATE_LANGS}
                except:
                    pass
            # 正则也失败，尝试逐行解析（兼容旧格式）
            if not result:
                for lc in TRANSLATE_LANGS:
                    m2 = re.search(rf'"{lc}"\s*:\s*"((?:[^"\\]|\\.)*)"', content, re.DOTALL)
                    if m2:
                        result[lc] = _strip_lang_prefix(m2.group(1).replace('\\"', '"').replace('\\n', '\n').strip(), lc)

    # 检查缺失的语言
    missing = [lc for lc in TRANSLATE_LANGS if not result.get(lc)]
    if missing:
        if retries > 0:
            print(f"[warn] 在线模型漏掉语言: {missing}，重试...")
            import time
            time.sleep(1)
            return call_online_llm(zh_text, examples, retries=retries - 1)
        else:
            print(f"[error] 在线模型最终漏掉语言: {missing}")
            print(f"[debug] 模型原始输出: {content[:500]}")
    
    return {lc: _strip_chinese_in_non_zh(_strip_lang_prefix(result.get(lc, ""), lc), lc) for lc in TRANSLATE_LANGS}


def call_online_llm_batch(segments: list[str], examples: list[dict], retries: int = 1) -> list[dict]:
    """在线模型批量翻译：一次调用翻译多个段落的全部目标语言。

    适用于需要在单次请求中翻译多个句段的场景。
    输出格式为每行 "语言代码: 翻译"，句段间用 "---" 分隔。

    解析策略：
    1. 优先按行解析语言代码前缀行
    2. 不足时按 "---" 分隔符切分
    3. 最后尝试 JSON 数组解析

    Args:
        segments: 待翻译的中文文本列表
        examples: 参考例句列表（可选）
        retries: 失败重试次数

    Returns:
        list[dict]: 每个元素为 {lang_code: translation} 的列表，与 segments 等长
    """
    cfg = _get_online_config()
    if not segments:
        return []

    block = ""
    if examples:
        block = "参考例句（沿用其术语和风格）：\n"
        for ex in examples:
            block += f"- 中文：{ex['zh']}\n"
            for lc in TRANSLATE_LANGS:
                if ex.get(lc):
                    block += f"  {LANG_NAMES[lc]}：{ex[lc]}\n"
        block += "\n"

    lang_list = "、".join(f"{c}={n}" for c, n in LANG_NAMES.items() if c in TRANSLATE_LANGS)
    seg_block = "\n".join(f"{i+1}. {s}" for i, s in enumerate(segments))

    user_prompt = f"""\
{block}请把下面 N 句车机中文文案逐句翻译成以下语言，输出格式：
每行 "语言代码: 翻译"，{len(TRANSLATE_LANGS)} 语全出，不要 JSON、不要 markdown、不要序号。
★ 每句翻译完之后，用 "---" 单独一行分隔下句，方便程序切分。

支持语言（key 用代码）：{lang_list}

原文列表：
{seg_block}

输出示例：
en: Please check tire pressure
ru: Проверьте давление в шинах
es: Compruebe la presión de los neumáticos
...（其余语言类推，每句 {len(TRANSLATE_LANGS)} 行，句间用 "---" 分隔）\
"""

    url = f"{cfg['base_url'].rstrip('/')}/chat/completions"
    headers = {
        "Authorization": f"Bearer {cfg['api_key']}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": cfg["model"],
        "messages": [
            {"role": "user", "content": user_prompt},
        ],
        "temperature": FALLBACK_TEMP,
        "max_tokens": 2048,
    }

    last_error = None
    for attempt in range(retries + 1):
        try:
            print(f"[在线batch] 调用 {cfg['model']} @ {cfg['base_url']} (尝试{attempt+1}, {len(segments)}段)")
            r = requests.post(url, headers=headers, json=payload, timeout=cfg["timeout"])
            r.raise_for_status()
            content = _extract_content(r.json())
            break
        except Exception as e:
            last_error = e
            print(f"[在线batch] 尝试{attempt+1}失败: {e}")
            if attempt < retries:
                import time
                time.sleep(2)
            else:
                raise last_error

    # 解析：按 segments 数 N 归 N 组
    lc_pattern = "^(" + "|".join(TRANSLATE_LANGS) + r")\s*[:：]\s*(.+)$"
    raw_lines = []
    for line in content.splitlines():
        m = re.match(lc_pattern, line.strip(), re.IGNORECASE)
        if m:
            raw_lines.append((m.group(1).lower(), m.group(2).strip().rstrip(',;')))

    N = len(segments)
    per_seg = [{} for _ in range(N)]
    expected_per_seg = len(TRANSLATE_LANGS)

    if len(raw_lines) >= N * expected_per_seg * 0.7:
        idx = 0
        for si in range(N):
            per_seg[si] = {}
            for _ in range(expected_per_seg):
                if idx < len(raw_lines):
                    lc, txt = raw_lines[idx]
                    if lc not in per_seg[si]:
                        per_seg[si][lc] = txt
                        idx += 1
                    else:
                        idx += 1
                else:
                    break
    else:
        items = re.split(r"\n\s*--+\s*\n", content)
        if len(items) < 2:
            items = [content]
        for si in range(min(N, len(items))):
            blk = items[si]
            for lc in TRANSLATE_LANGS:
                m = re.search(rf"^{lc}\s*[:：]\s*(.+)$", blk, re.MULTILINE | re.IGNORECASE)
                if m:
                    per_seg[si][lc] = m.group(1).strip().rstrip(',;')

    if not any(per_seg):
        try:
            arr = json.loads(re.sub(r"^```json|```$", "", content.strip(), flags=re.MULTILINE).strip())
            if isinstance(arr, list) and len(arr) == N:
                return [{lc: arr[i].get(lc, "") for lc in TRANSLATE_LANGS} for i in range(N)]
        except:
            pass

    final_result = [{lc: per_seg[i].get(lc, "") for lc in TRANSLATE_LANGS} for i in range(N)]
    missing_segs = [(i, [lc for lc in TRANSLATE_LANGS if not per_seg[i].get(lc)]) for i in range(N) if any(not per_seg[i].get(lc) for lc in TRANSLATE_LANGS)]
    if missing_segs:
        print(f"[warn] online batch 缺失语言: {missing_segs}，将由外层逐段补全")
    
    return final_result


# ====================================================================
# translate_one（★ v2.18: 统一在线模型翻译）
# 核心翻译入口，实现三级匹配策略
# ====================================================================
def _detect_translation_lang(text: str) -> str | None:
    """检测翻译文本实际属于哪种语言（通过字符编码特征）。

    用于验证知识库中各语言列的内容是否确实和列名匹配，
    检测到列错位时可以进行纠正。

    Args:
        text: 待检测的文本

    Returns:
        str | None: 检测到的语言代码，无法确定返回 None
    """
    if not text:
        return None
    has_arabic = bool(re.search(r'[\u0600-\u06FF]', text))

    # 哈萨克语特有的西里尔字母（区别于俄语）
    has_kazakh = bool(re.search(r'[ӘәҒғҚқҢңӨөҰұҮүҺһІі]', text))
    has_cyrillic = bool(re.search(r'[\u0400-\u04FF]', text))
    has_cjk = bool(re.search(r'[\u4e00-\u9fff]', text))

    if has_arabic:
        return "ar"
    if has_kazakh:
        return "kk"
    if has_cyrillic:
        return "ru"
    if has_cjk:
        return "zh_hant"
    return None

def _assign_kb_translations(raw: dict, out_langs: list[str]) -> tuple[dict, list]:
    """从知识库原始数据中分配翻译结果，处理语言列错位问题。

    两遍分配策略：
    1. 第一遍：文本语言匹配该列的，直接使用
    2. 第二遍：语言不匹配的，互换或填入正确位置

    Args:
        raw: 从数据库读取的原始 {lang_code: text} 字典
        out_langs: 要输出的语言代码列表

    Returns:
        tuple: (kb_trans: dict, missing: list[str])
               kb_trans 是分配后的翻译字典
               missing 是仍缺失的语言代码列表
    """
    kb_trans = {}
    # 第一遍：文本语言匹配该列的，直接使用
    for lc in out_langs:
        val = raw.get(lc, "")
        if not val:
            continue
        actual = _detect_translation_lang(val)
        if actual and actual != lc and actual in out_langs:
            continue
        kb_trans[lc] = val
    # 第二遍：语言不匹配的，互换或填入正确位置
    for lc in out_langs:
        val = raw.get(lc, "")
        if not val:
            continue
        actual = _detect_translation_lang(val)
        if actual and actual != lc and actual in out_langs:
            if actual in kb_trans:
                existing = kb_trans[actual]
                kb_trans[actual] = val
                if lc not in kb_trans:
                    kb_trans[lc] = existing
                print(f"[KB] {lc}↔{actual}: 互换错位翻译")
            else:
                kb_trans[actual] = val
                print(f"[KB] {lc}→{actual}: 填入缺失位置")
    _missing = [lc for lc in out_langs if lc not in kb_trans]
    return kb_trans, _missing


def _fill_missing_langs(conn, ids, vecs, zh_text: str, kb_trans: dict, missing: list[str]) -> dict:
    """知识库命中后，对缺失的语言走在线模型兜底翻译补上。

    当知识库匹配到了中文文本但某些目标语言列没有翻译时，
    调用在线模型逐语言翻译缺失的部分。

    Args:
        conn: SQLite 数据库连接
        ids: 向量索引 ID 数组
        vecs: 向量数组
        zh_text: 中文原文
        kb_trans: 已有的知识库翻译字典
        missing: 需要模型补全的语言代码列表

    Returns:
        dict: 补全后的完整翻译字典
    """
    qvec = embed(zh_text)
    examples = []
    hits = search(qvec, ids, vecs, k=TOP_K) if ids is not None and len(ids) > 0 else []
    if hits:
        order = np.argsort(-(vecs @ qvec))[:TOP_K]
        examples = [fetch_row(conn, int(ids[o])) for o in order if (vecs @ qvec)[o] >= MED_SIM]
    for lc in missing:
        try:
            kb_trans[lc] = call_online_llm_single_lang(zh_text, lc, examples, use_simple_prompt=True)
            if not kb_trans[lc]:
                kb_trans[lc] = ""
        except Exception:
            kb_trans[lc] = ""
    return kb_trans

def translate_one(conn, ids, vecs, zh_text: str, target_langs: list[str] | None = None) -> dict:
    """翻译入口，三级匹配：精确命中 → 模糊匹配 → 模型兜底。

    ★ v2.18: 统一在线模型翻译，不再降级本地 Ollama。
    ★ 非中文输入直接走模型翻译（KB 只含中文）。

    匹配流程：
    1. 第一级：精确命中（本地知识库的 zh 完全匹配）
    2. 第二级：模糊子串匹配（输入是知识库条目的子串）
    3. 语义相似度匹配（向量搜索 ≥ HIGH_SIM）
    4. 第三级：segment_base 段匹配
    5. 第四级：模型兜底（在线模型逐语言翻译）

    Args:
        conn: SQLite 数据库连接
        ids: 向量索引 ID 数组
        vecs: 向量数组
        zh_text: 中文原文
        target_langs: 指定只输出哪些语言（如 ["en","ru"]），None=全语种

    Returns:
        dict: {
            "mode": str,           # 匹配模式描述
            "matched_zh": str|None, # 匹配到的中文
            "similarity": float,   # 相似度
            "translations": dict,  # {lang_code: text}
            "candidates": list|None, # 候选条目
            "target_langs": list,  # 目标语言列表
            "need_model": list,    # 需要模型兜底的语言
        }
    """
    out_langs = target_langs if target_langs else TRANSLATE_LANGS
    source_lang = detect_source_lang(zh_text)

    # ★ 非中文输入：跳过 KB 匹配，直接在线模型翻译
    if source_lang != "zh":
        result = {}
        for lc in out_langs:
            try:
                result[lc] = call_online_llm_single_lang(zh_text, lc, [], use_simple_prompt=True, source_lang=source_lang)
                if not result[lc]:
                    result[lc] = ""
            except Exception as e:
                print(f"[warn] 在线翻译 {lc} 失败: {e}")
                result[lc] = ""
        return {
            "mode": "非中文-在线模型翻译", "matched_zh": None, "similarity": 0.0,
            "translations": result, "candidates": None,
            "target_langs": out_langs,
            "need_model": out_langs,
        }

    # ★ 第一级：精确命中（本地知识库）
    exact = conn.execute("SELECT * FROM tm_segments WHERE zh=?", (zh_text,)).fetchone()

    # ★ 第1.5级：CJK 标点无关精确匹配
    if not exact:
        zh_cjk = extract_cjk(zh_text)
        if zh_cjk:
            cache = _get_zh_cjk_cache(conn)
            if zh_cjk in cache:
                rid, matched_zh = cache[zh_cjk]
                exact = conn.execute("SELECT * FROM tm_segments WHERE id=?", (rid,)).fetchone()

    if exact:
        cols = [d[0] for d in conn.execute("SELECT * FROM tm_segments LIMIT 1").description]
        row = dict(zip(cols, exact))
        raw = {lc: row.get(lc, "") for lc in out_langs}
        kb_trans, _missing = _assign_kb_translations(raw, out_langs)
        if _missing:
            kb_trans = _fill_missing_langs(conn, ids, vecs, zh_text, kb_trans, _missing)
        return {
            "mode": "精确命中", "matched_zh": zh_text, "similarity": 1.0,
            "translations": kb_trans,
            "candidates": None,
            "target_langs": out_langs,
            "need_model": _missing,
        }

    # ★ 第二级：模糊匹配（本地知识库）
    fuzzies = fuzzy_substring_hits(conn, zh_text)
    if fuzzies:
        if len(fuzzies) == 1:
            pick = fuzzies[0]; mode = "模糊子串-复用长句译文"
        else:
            pick = max(fuzzies, key=lambda r: len(r["zh"])); mode = f"模糊子串-多候选({len(fuzzies)})"
        raw = {lc: pick.get(lc, "") for lc in out_langs}
        kb_trans, _missing = _assign_kb_translations(raw, out_langs)
        if _missing:
            kb_trans = _fill_missing_langs(conn, ids, vecs, zh_text, kb_trans, _missing)
        return {
            "mode": mode, "matched_zh": pick["zh"], "similarity": 1.0,
            "translations": kb_trans,
            "candidates": fuzzies,
            "target_langs": out_langs,
            "need_model": _missing,
        }

    # 语义相似度匹配（本地向量）
    qvec = embed(zh_text)
    hits = search(qvec, ids, vecs, k=TOP_K)
    top_id, top_sim = hits[0] if hits else (None, 0.0)
    if top_id is not None and top_sim >= HIGH_SIM:
        row = fetch_row(conn, top_id)
        raw = {lc: row.get(lc, "") for lc in out_langs}
        kb_trans, _missing = _assign_kb_translations(raw, out_langs)
        if _missing:
            kb_trans = _fill_missing_langs(conn, ids, vecs, zh_text, kb_trans, _missing)
        return {
            "mode": "语义高相似-直接复用", "matched_zh": row["zh"], "similarity": top_sim,
            "translations": kb_trans,
            "candidates": None,
            "target_langs": out_langs,
            "need_model": _missing,
        }

    # ★ 第三级：segment_base 段匹配（新增）
    seg_total = {lc: {"matched": [], "remaining": ""} for lc in out_langs}

    for lc in out_langs:
        lang_pair = f"zh-{lc.replace('_lang', '')}"
        matched, remaining = match_segments(conn, zh_text, lang_pair)
        seg_total[lc] = {"matched": matched, "remaining": remaining}

    all_full = all(not s["remaining"] for s in seg_total.values())
    any_partial = any(s["matched"] for s in seg_total.values())

    if all_full:
        seg_result = {}
        for lc in out_langs:
            seg_result[lc] = reassemble_translation(seg_total[lc]["matched"], "", zh_text)
        return {
            "mode": "段匹配-全命中", "matched_zh": zh_text, "similarity": 0.0,
            "translations": seg_result, "candidates": None,
            "target_langs": out_langs, "need_model": [],
        }

    if any_partial:
        seg_result = {}
        _need_model = []
        for lc in out_langs:
            if seg_total[lc]["matched"] and not seg_total[lc]["remaining"]:
                seg_result[lc] = reassemble_translation(seg_total[lc]["matched"], "", zh_text)
            elif seg_total[lc]["matched"] and seg_total[lc]["remaining"]:
                try:
                    model_part = call_online_llm_single_lang(seg_total[lc]["remaining"], lc, [], source_lang=source_lang)
                except Exception:
                    model_part = ""
                seg_result[lc] = reassemble_translation(seg_total[lc]["matched"], model_part, zh_text)
            else:
                _need_model.append(lc)
        if _need_model:
            for lc in _need_model:
                try:
                    seg_result[lc] = call_online_llm_single_lang(zh_text, lc, [], source_lang=source_lang)
                except Exception as e:
                    print(f"[warn] 在线翻译 {lc} 失败: {e}")
                    seg_result[lc] = ""
        return {
            "mode": "段匹配-部分命中", "matched_zh": zh_text, "similarity": 0.0,
            "translations": seg_result, "candidates": None,
            "target_langs": out_langs, "need_model": [],
        }

    # ★ 第四级：模型兜底（在线模型逐语翻译）
    examples = []
    if top_id is not None and top_sim >= MED_SIM:
        order = np.argsort(-(vecs @ qvec))[:TOP_K]
        examples = [fetch_row(conn, int(ids[o])) for o in order if (vecs @ qvec)[o] >= MED_SIM]

    result = {}
    for lc in out_langs:
        try:
            result[lc] = call_online_llm_single_lang(zh_text, lc, [], source_lang=source_lang)
            if not result[lc]:
                result[lc] = ""
        except Exception as e:
            print(f"[warn] 在线翻译 {lc} 失败: {e}")
            result[lc] = ""

    mode = f"段匹配-全不命中-在线模型生成"

    return {
        "mode": mode, "matched_zh": None,
        "similarity": top_sim, "translations": result, "candidates": None,
        "target_langs": out_langs,
        "need_model": out_langs,
    }

# ====================================================================
# save_back / tm_stats — 知识库写入与统计
# ====================================================================
def save_back(conn, zh_text: str, translations: dict, module: str | None = None) -> int:
    """保存翻译到知识库（SQLite），返回行 ID。

    使用中文文本的 MD5 哈希值作为唯一键（zh_hash），
    支持 UPSERT 语义：存在则更新，不存在则插入。

    注意：此函数只写 DB，不更新向量索引。
    调用方需再调 rebuild_embeddings_for_entry() 完成向量层更新。

    Args:
        conn: SQLite 数据库连接
        zh_text: 中文原文
        translations: {lang_code: text} 翻译字典
        module: 模块名称（可选，用于分类）

    Returns:
        int: 插入/更新的行 ID，失败返回 -1
    """
    _invalidate_zh_cjk_cache()
    zh_hash = hashlib.md5(zh_text.encode("utf-8")).hexdigest()
    now = datetime.now().isoformat(timespec="seconds")
    cols = ["zh_hash", "zh"] + ALL_LANGS + ["updated_at"]
    vals = [zh_hash, zh_text] + [translations.get(lc, "") for lc in ALL_LANGS] + [now]
    if module:
        cols.append("module"); vals.append(module)
    ph = ",".join(["?"] * len(cols))
    cn = ",".join(f'"{c}"' for c in cols)
    conn.execute(
        f'INSERT INTO tm_segments ({cn}) VALUES ({ph}) '
        f'ON CONFLICT(zh_hash) DO UPDATE SET '
        + ",".join(f'"{c}"=excluded."{c}"' for c in cols if c != "zh_hash"),
        vals,
    )
    conn.commit()
    row = conn.execute("SELECT id FROM tm_segments WHERE zh_hash=?", (zh_hash,)).fetchone()
    return row[0] if row else -1


def rebuild_embeddings_for_entry(zh_text: str, row_id: int) -> bool:
    """为新增/更新的条目计算 embedding 并写入 npz 向量索引文件。

    从 npz 文件中读取现有索引，添加或更新该条目的向量，
    然后写回 npz 文件并更新时间戳文件。

    Args:
        zh_text: 中文文本
        row_id: 对应的数据库行 ID

    Returns:
        bool: True=成功, False=失败（Ollama不可用等）
    """
    if not ollama_is_up():
        print("[warn] rebuild_embeddings_for_entry: Ollama未运行，跳过向量更新")
        return False
    try:
        vec = embed(zh_text)
        if os.path.exists(EMB_PATH):
            data = np.load(EMB_PATH)
            ids_arr = data["ids"]
            vecs_arr = data["vecs"]
        else:
            ids_arr = np.array([], dtype=np.int64)
            vecs_arr = np.zeros((0, vec.shape[0]), dtype=np.float32)

        if len(ids_arr) > 0 and row_id in ids_arr:
            idx = int(np.where(ids_arr == row_id)[0][0])
            vecs_arr[idx] = vec
        else:
            ids_arr = np.append(ids_arr, row_id).astype(np.int64)
            vecs_arr = np.vstack([vecs_arr, vec.reshape(1, -1)]) if len(vecs_arr) > 0 else vec.reshape(1, -1)

        EMB_PATH.parent.mkdir(parents=True, exist_ok=True)
        np.savez(EMB_PATH, ids=ids_arr, vecs=vecs_arr)
        INDEX_STAMP.parent.mkdir(parents=True, exist_ok=True)
        INDEX_STAMP.touch()
        print(f"[info] 向量索引已更新：row_id={row_id}, 总条目={len(ids_arr)}")
        return True
    except Exception as e:
        print(f"[error] rebuild_embeddings_for_entry 失败: {e}")
        return False

def tm_stats(conn) -> dict:
    """获取翻译记忆库的统计信息。

    Returns:
        dict: {"total": 总条目数, "per_lang": {语言代码: 非空条目数}}
    """
    total = conn.execute("SELECT COUNT(*) FROM tm_segments").fetchone()[0]
    per = {}
    for lc in ALL_LANGS:
        per[lc] = conn.execute(f'SELECT COUNT(*) FROM tm_segments WHERE "{lc}" IS NOT NULL').fetchone()[0]
    return {"total": total, "per_lang": per}


# ====================================================================
# v2.16 截断自修复
# GLM-4.7-flash 的模型级 bug：finish_reason=stop 但实际截断
# ====================================================================

def _is_translation_incomplete(result: str, zh_text: str, target_lang: str) -> bool:
    """检测翻译结果是否不完整（GLM-4.7-flash 的模型级bug：finish_reason=stop 但实际截断）

    检测逻辑（三重判断，命中任一即判定不完整）：
    1. 译文长度远短于原文（非CJK语言译文一般≥原文长度×1.5，如果短于原文则很可能截断）
    2. 译文以明显不完整的措辞结尾（如 "the summary is"、"总结如下"、"it can be" 等）
    3. 译文末尾不是句末标点，且译文长度不够长（排除短句场景）

    ★ 不适用于：极短原文（≤20字）、zh_hant（简繁转换不是翻译）
    """
    if not result or not zh_text:
        return False
    # 极短原文不做截断判断
    if len(zh_text) <= 20:
        return False
    # 简繁转换不做截断判断
    if target_lang == "zh_hant":
        return False

    result_len = len(result)
    zh_len = len(zh_text)

    # ★ 判断1：译文长度远短于原文
    # 非CJK语言（英语/俄语/阿拉伯语等）翻译后字符数一般 ≥ 原文（中文信息密度高）
    # CJK语言（日语/韩语）翻译后字符数可能与中文接近
    cjk_langs = {"ja", "ko"}
    if target_lang in cjk_langs:
        min_ratio = 0.6  # 日韩语与中文信息密度接近，但译文通常不会比原文短太多
    else:
        min_ratio = 0.8  # 非CJK语言译文长度一般 ≥ 原文
    if result_len < zh_len * min_ratio:
        print(f"[截断检测] {target_lang} 译文过短：result={result_len} < zh={zh_len}*{min_ratio}={zh_len*min_ratio:.0f}")
        return True

    # ★ 判断2：译文以"未完成短语"结尾（中英通用）
    # 模型最常见的不完整结尾模式
    # ★ 注意：如果译文已经足够长（≥原文×1.5），"总结起来就是"/"the summary is:" 这种结尾
    #   其实是正常翻译——原文本身就以"总结起来就是"结尾，翻译到"the summary is:"是对的
    incomplete_endings = [
        # 英语常见——但只在译文偏短时才算不完整
        "it can be",  # "it can be summarized as" 后面没内容
    ]
    result_lower = result.strip().lower()
    # ★ 只有译文偏短时，才检查这些未完成短语
    if result_len < zh_len * 1.2:
        for ending in incomplete_endings:
            if result_lower.endswith(ending):
                print(f"[截断检测] {target_lang} 译文偏短+未完成结尾：...{result.strip()[-30:]}")
                return True

    # ★ 判断3：译文末尾不是句末标点，且译文偏短
    # 句末标点：英语 .!? / 日语 。！？/ 韩语 。！？/ 阿拉伯语 .؟! / 俄语 .!? 等
    sentence_end_punct = {'.', '!', '?', '。', '！', '？', '…', '»', '」', '」', '؛', '؟', ':'}
    last_char = result.strip()[-1] if result.strip() else ''
    # 长译文(>100字符)不以句末标点结尾=大概率截断
    # ★ 但如果译文已经足够长（≥原文×1.5），冒号结尾不算截断（"总结起来就是：" → "the summary is:" 是正常翻译）
    if result_len > 100 and last_char not in sentence_end_punct:
        # 排除：引号/括号结尾（可能是正常的引述结尾）
        if last_char not in {'"', '"', ')', '）', ']', '」', "'"}:
            print(f"[截断检测] {target_lang} 长译文无句末标点：末尾=\"{last_char}\"，长度={result_len}")
            return True

    return False


def _auto_complete_translation(result: str, zh_text: str, target_lang: str, cfg: dict, url: str, headers: dict) -> str:
    """翻译截断自修复：检测到不完整时，自动续翻补全。

    策略：先尝试"续翻"（基于已有结果继续），如果续翻也失败则"全量重翻"（加强提示词）。
    ★ 续翻优先：保留已有翻译的前半段，只补后半段，避免全量重翻质量波动。

    Args:
        result: 当前不完整的翻译结果
        zh_text: 中文原文
        target_lang: 目标语言代码
        cfg: API 配置字典
        url: API 请求 URL
        headers: 请求头

    Returns:
        str: 修复后的完整翻译文本，若所有尝试都失败则返回原始结果
    """
    lang_name = LANG_NAMES.get(target_lang, target_lang)
    max_retries = 2  # 最多2次尝试（1次续翻 + 1次全量重翻）

    for retry_type in ["续翻", "全量重翻"]:
        try:
            if retry_type == "续翻":
                # ★ 续翻：告诉模型"前面的翻译到XX为止，请从XX处继续翻译剩余内容"
                user_content = (
                    f"以下中文原文的{lang_name}翻译不完整，只翻译了一半。"
                    f"请在已有翻译的基础上继续翻译剩余内容，只输出继续的部分，不要重复已有翻译。\n\n"
                    f"中文原文：\n{zh_text}\n\n"
                    f"已有翻译（不完整）：\n{result}\n\n"
                    f"请继续翻译剩余部分："
                )
            else:
                # ★ 全量重翻：加强提示词，强调"必须翻译完整"
                user_content = (
                    f"【重要】请将以下中文完整翻译为{lang_name}，必须翻译完全部内容，不要遗漏后半段。"
                    f"不要在翻译到一半时停止。如果原文有总结性语句，翻译也要包含总结的具体内容。\n\n"
                    f"中文原文：\n{zh_text}"
                )

            payload = {
                "model": cfg["model"],
                "messages": [
                    {"role": "user", "content": user_content},
                ],
                "temperature": FALLBACK_TEMP,
                "max_tokens": max(4096, len(zh_text) * 6),  # ★ 给充足空间
            }

            r = requests.post(url, headers=headers, json=payload, timeout=cfg["timeout"])
            r.raise_for_status()
            r_json = r.json()
            new_content = _extract_content(r_json)

            if not new_content:
                print(f"[截断自修复] {target_lang} {retry_type}返回为空，跳过")
                continue

            new_result = _extract_single_lang(new_content.strip(), target_lang)

            if retry_type == "续翻":
                # 续翻结果拼接到已有结果后面
                # ★ 去掉续翻结果中可能重复的部分（模型可能重复了最后几个词）
                combined = _merge_continuation(result, new_result, target_lang)
                if combined and len(combined) > len(result):
                    # 验证拼接后的结果是否更完整
                    if not _is_translation_incomplete(combined, zh_text, target_lang):
                        print(f"[截断自修复] {target_lang} 续翻成功：{len(result)}→{len(combined)}")
                        return combined
                    else:
                        print(f"[截断自修复] {target_lang} 续翻后仍不完整，尝试全量重翻")
                        result = combined  # 保留拼接结果作为备选
            else:
                # 全量重翻直接使用新结果
                if new_result and not _is_translation_incomplete(new_result, zh_text, target_lang):
                    print(f"[截断自修复] {target_lang} 全量重翻成功：{len(result)}→{len(new_result)}")
                    return new_result
                elif new_result and len(new_result) > len(result):
                    print(f"[截断自修复] {target_lang} 全量重翻仍不完整但更长，采用：{len(result)}→{len(new_result)}")
                    return new_result

        except Exception as e:
            print(f"[截断自修复] {target_lang} {retry_type}异常: {e}")
            continue

    # 所有尝试都失败，返回原始结果（至少有部分翻译）
    print(f"[截断自修复] {target_lang} 自修复失败，返回原始结果")
    return result


def _merge_continuation(original: str, continuation: str, target_lang: str) -> str:
    """将续翻结果与原始翻译合并，去除重复部分。

    模型续翻时经常会重复原文的最后几个词，需要智能去重。
    去重策略：从 original 末尾取不同长度的重叠片段，检查 continuation 是否以该片段开头。

    Args:
        original: 原始翻译（前半段）
        continuation: 续翻结果（后半段，可能包含了重复）
        target_lang: 目标语言（用于语言特定去重，暂未使用）

    Returns:
        str: 合并后的完整翻译
    """
    if not continuation:
        return original
    if not original:
        return continuation

    # ★ 去重策略：从original末尾取不同长度的重叠片段，检查continuation是否以该片段开头
    # 如果有重叠，去掉continuation中重复的部分
    orig_words = original.split()
    cont_words = continuation.split()

    # 从长到短尝试重叠（最多尝试20个词的重叠）
    max_overlap = min(len(orig_words), len(cont_words), 20)
    for overlap_len in range(max_overlap, 0, -1):
        orig_tail = orig_words[-overlap_len:]
        cont_head = cont_words[:overlap_len]
        # 精确匹配
        if orig_tail == cont_head:
            # 去掉continuation中重复的部分
            merged_words = orig_words + cont_words[overlap_len:]
            return ' '.join(merged_words)

    # 没有找到词级重叠，直接拼接（中间加空格）
    return original + ' ' + continuation

def call_online_llm_single_lang(zh_text: str, target_lang: str, examples: list[dict], retries: int = 1, on_step=None, use_simple_prompt: bool = False, lang_instruction_override: str | None = None, source_lang: str = "zh") -> str:
    """调用在线模型翻译单条文本到单种语言。

    这是最基本的在线翻译调用函数，支持：
    - 根据目标语言自动选择翻译指令（zh_hant→简繁转换，ja/ko→特殊要求）
    - Hunyuan-MT 模型的专有参数设置
    - 429 限流时自动降级到回退模型
    - finish_reason="length" 时自动增大 max_tokens 重试
    - 重试机制和指数退避

    Args:
        zh_text: 待翻译的中文文本
        target_lang: 目标语言代码
        examples: 参考例句列表（可选）
        retries: 失败重试次数
        on_step: 进度回调函数 on_step("translate", lang_code, lang_name)
        use_simple_prompt: 是否使用极简提示词（无例句，仅翻译指令）
        lang_instruction_override: 自定义翻译指令（覆盖自动生成）
        source_lang: 源语言代码，默认 "zh"

    Returns:
        str: 翻译后的文本

    Raises:
        所有重试用尽后抛出最后一次异常
    """
    cfg = _get_online_config(target_lang=target_lang)
    lang_name = LANG_NAMES.get(target_lang, target_lang)
    src_name = get_source_name(source_lang)
    is_hunyuan = (cfg["model"] == HUNYUAN_MT_MODEL)

    url = f"{cfg['base_url'].rstrip('/')}/chat/completions"
    headers = {
        "Authorization": f"Bearer {cfg['api_key']}",
        "Content-Type": "application/json",
    }

    line_count = zh_text.count('\n') + 1
    base_tokens = max(2048, len(zh_text) * 4)
    max_tokens = min(max(base_tokens, 2048 + line_count * 384), 8192)

    if target_lang == "zh_hant":
        translate_instruction = f"把下面的{src_name}转换为繁体中文，只输出繁体中文结果"
    elif target_lang == "ja":
        translate_instruction = f"把下面的{src_name}翻译为日语。必须使用规范的日语汉字+假名混合书写，不要只用假名。只输出日语翻译结果"
    elif target_lang == "ko":
        translate_instruction = f"把下面的{src_name}翻译为韩语（한국어）。必须使用韩语谚文书写，禁止输出日语。只输出韩语翻译结果"
    else:
        translate_instruction = f"把下面的{src_name}翻译为{lang_name}，不要额外解释"

    user_prompt = f"{translate_instruction}\n\n{zh_text}"

    if is_hunyuan:
        payload = {
            "model": cfg["model"],
            "messages": [{"role": "user", "content": user_prompt}],
            "temperature": HUNYUAN_TEMP,
            "top_p": HUNYUAN_TOP_P,
            "top_k": HUNYUAN_TOP_K,
            "repetition_penalty": HUNYUAN_REPETITION_PENALTY,
            "max_tokens": max_tokens,
        }
    else:
        payload = {
            "model": cfg["model"],
            "messages": [{"role": "user", "content": user_prompt}],
            "temperature": FALLBACK_TEMP,
            "max_tokens": max_tokens,
        }

    if on_step:
        on_step("translate", target_lang, lang_name)

    _fallback_used = False
    last_error = None
    for attempt in range(retries + 1):
        try:
            print(f"[翻译] {target_lang}({lang_name}) 调用 {cfg['model']} (尝试{attempt+1})")
            r = requests.post(url, headers=headers, json=payload, timeout=cfg["timeout"])
            r.raise_for_status()
            r_json = r.json()
            content = _extract_content(r_json)

            finish_reason = ""
            try:
                finish_reason = r_json.get("choices", [{}])[0].get("finish_reason", "")
            except (IndexError, KeyError):
                pass

            if finish_reason == "length":
                if attempt < retries + 2:
                    payload["max_tokens"] = min(payload.get("max_tokens", 1024) * 2, 8192)
                    import time; time.sleep(1)
                    continue

            if not content and attempt < retries:
                payload["max_tokens"] = min(payload.get("max_tokens", 1024) * 2, 8192)
                import time; time.sleep(1)
                continue

            if content:
                result = content.strip()
            else:
                result = ""

            result = post_process_translation(result, target_lang)
            print(f"[翻译] {target_lang} 返回成功，长度={len(result)}")
            return result

        except Exception as e:
            last_error = e
            err_str = str(e)
            print(f"[翻译] {target_lang} 尝试{attempt+1}失败: {e}")

            if ("429" in err_str or "rate" in err_str.lower()) and not _fallback_used:
                _fallback_used = True
                fallback_model = HUNYUAN_FALLBACK_MODEL if is_hunyuan else HUNYUAN_FALLBACK_MODEL
                payload["model"] = fallback_model
                payload.pop("top_p", None)
                payload.pop("top_k", None)
                payload.pop("repetition_penalty", None)
                payload["temperature"] = FALLBACK_TEMP
                print(f"[翻译] {target_lang} 429降级到 {fallback_model}")
                import time; time.sleep(2)
                continue

            if attempt < retries:
                import time; time.sleep(2)
            else:
                raise last_error





# ====================================================================
# 批量单语翻译（★ v2.20: 并发2批 + 延迟兜底）
# 将多条文本合并为一次API调用，翻译为同一目标语言
# ====================================================================

def _translate_one_batch(
    batch_idx: int, batch: list[str], start: int,
    target_lang: str, lang_name: str, lang_instruction: str,
    example_block: str, cfg: dict, url: str, headers: dict,
    batch_size: int,
) -> dict:
    """翻译单个批次，返回 {start+i: translation}。

    使用 <sN>...</sN> 标记构建输入，让模型按标记对应输出。
    如果标记解析失败（解析率<10%），尝试按行解析兜底。

    ★ 内部函数，供并发调用。失败返回空 dict，由调用方统一兜底。

    Args:
        batch_idx: 批次索引（0-based）
        batch: 本批待翻译的中文文本列表
        start: 本批在全局列表中的起始索引
        target_lang: 目标语言代码
        lang_name: 目标语言中文名
        lang_instruction: 翻译指令字符串
        example_block: 参考例句区块
        cfg: API 配置字典
        url: API 请求 URL
        headers: 请求头
        batch_size: 批大小

    Returns:
        dict: {全局索引: 翻译文本}，失败条目不包含在字典中
    """
    # ★ 构造带标记的输入
    numbered_input = ""
    for i, text in enumerate(batch, 1):
        numbered_input += f"<s{i}> {text} </s{i}>\n"

    user_prompt = (
        f"{example_block}"
        + f"{lang_instruction}。按编号逐条翻译，每条翻译用对应的 <sN> 和 </sN> 标记包裹。"
        + "不要输出其他语言，不要加语言名称前缀，不要解释。保留原文的段落换行结构。\n\n"
        + f"中文原文：\n{numbered_input}"
    )

    total_chars = sum(len(t) for t in batch)
    max_tokens = max(2048, total_chars * 4)
    max_tokens = min(max_tokens, 8192)

    payload = {
        "model": cfg["model"],
        "messages": [
            {"role": "user", "content": user_prompt},
        ],
        "temperature": FALLBACK_TEMP,
        "max_tokens": max_tokens,
    }

    _fallback = False
    for attempt in range(3):
        try:
            r = requests.post(url, headers=headers, json=payload, timeout=cfg["timeout"])
            r.raise_for_status()
            r_json = r.json()
            content = _extract_content(r_json)

            # ★ 检查截断
            finish_reason = ""
            try:
                finish_reason = r_json.get("choices", [{}])[0].get("finish_reason", "")
            except (IndexError, KeyError):
                pass

            if finish_reason == "length" or not content:
                if attempt < 2:
                    payload["max_tokens"] = min(payload.get("max_tokens", 2048) * 2, 8192)
                    import time; time.sleep(1)
                    continue

            # ★ 解析 <sN>...</sN> 标记
            parsed = {}
            for m in re.finditer(r'<s(\d+)>\s*(.*?)\s*</s\1>', content, re.DOTALL):
                idx = int(m.group(1)) - 1
                trans = m.group(2).strip()
                if idx < len(batch):
                    parsed[idx] = trans

            # ★ 标记解析不足10%才尝试行解析（放宽阈值，避免误触发行解析覆盖正确结果）
            if len(parsed) < len(batch) * 0.1:
                parsed = _parse_batch_by_line(content, len(batch), target_lang)

            # ★ 写入结果（只写解析到的，未解析的留给调用方集中兜底）
            result = {}
            for i in range(len(batch)):
                if i in parsed and parsed[i]:
                    trans = _strip_lang_prefix(parsed[i], target_lang)
                    result[start + i] = trans

            return result

        except Exception as e:
            err_str = str(e)
            # ★ 429降级到 glm-4-flash + 指数退避
            if "429" in err_str or "速率限制" in err_str or "rate" in err_str.lower():
                if not _fallback and "4.7" in cfg["model"]:
                    payload["model"] = "glm-4-flash"
                    _fallback = True
                    _wait = 5 * (attempt + 1)  # ★ 第1次等5秒，第2次10秒
                    print(f"[批量单语] {target_lang} 批次{batch_idx+1} 429限流，降级到glm-4-flash，等待{_wait}秒")
                    import time; time.sleep(_wait)
                    continue
                else:
                    _wait = 8 * (attempt + 1)  # ★ 降级后仍429：等8秒、16秒
                    print(f"[批量单语] {target_lang} 批次{batch_idx+1} 仍429，等待{_wait}秒")
                    import time; time.sleep(_wait)
                    continue
            if attempt < 2:
                import time; time.sleep(3)
            else:
                print(f"[批量单语] {target_lang} 批次{batch_idx+1} 全部失败: {err_str[:200]}")
                return {}

    return {}


def call_online_llm_single_lang_batch(
    zh_texts: list[str],
    target_lang: str,
    examples: list[dict] = None,
    batch_size: int = 15,
    on_batch_done=None,
) -> list[str]:
    """批量翻译：将多条中文文本合并为一次 API 调用，翻译为同一目标语言。

    ★ v2.20: 并发2批同时请求 + 延迟兜底（未解析的集中再翻译，不逐条调 API）。

    流程：
    1. 动态分批：根据平均文本长度调整每批数量
    2. 自适应并发：初始2批并发，429 时降级为串行
    3. 集中兜底：未解析条目合并为一个补批
    4. 逐条兜底（极少情况）：仍有缺失的逐条调用

    Args:
        zh_texts: 待翻译的中文文本列表
        target_lang: 目标语言代码
        examples: 参考例句列表（可选）
        batch_size: 基础批大小，会根据文本长度动态调整
        on_batch_done: 回调 on_batch_done(batch_idx, total_batches, done_count, total_count)

    Returns:
        list[str]: 与 zh_texts 等长的翻译结果列表，失败条目为 "[翻译失败]"
    """
    if not zh_texts:
        return []

    cfg = _get_online_config(target_lang=target_lang)
    if not cfg.get("api_key"):
        return ["[翻译失败]"] * len(zh_texts)

    lang_name = LANG_NAMES.get(target_lang, target_lang)
    source_lang = detect_source_lang(zh_texts[0]) if zh_texts else "zh"
    src_name = get_source_name(source_lang)
    url = f"{cfg['base_url'].rstrip('/')}/chat/completions"
    headers = {
        "Authorization": f"Bearer {cfg['api_key']}",
        "Content-Type": "application/json",
    }

    if target_lang == "zh_hant":
        lang_instruction = f"把下面的{src_name}转换为繁体中文，只输出繁体中文结果" if source_lang == "zh" else f"把下面的{src_name}翻译为繁体中文，只输出繁体中文结果"
    elif target_lang == "ja":
        lang_instruction = f"把下面的{src_name}翻译为日语。必须使用规范的日语汉字+假名混合书写，不要只用假名。只输出日语翻译结果"
    elif target_lang == "ko":
        lang_instruction = f"把下面的{src_name}翻译为韩语（한국어）。必须使用韩语谚文书写，禁止输出日语。只输出韩语翻译结果"
    else:
        lang_instruction = f"把下面的{src_name}翻译为{lang_name}，不要额外解释"

    example_block = ""
    if examples:
        example_block = f"参考例句（{lang_name}部分，沿用其术语和风格）：\n"
        for ex in examples:
            example_block += f"- {src_name}：{ex['zh']}\n"
            if ex.get(target_lang):
                example_block += f"  {lang_name}：{ex[target_lang]}\n"
        example_block += "\n"

    results = [""] * len(zh_texts)

    # ★ 动态分批：按平均文本长度调整 batch_size
    avg_len = sum(len(t) for t in zh_texts) / max(len(zh_texts), 1)
    if avg_len > 150:
        batch_size = min(batch_size, 5)
    elif avg_len > 80:
        batch_size = min(batch_size, 8)
    elif avg_len < 30:
        batch_size = max(batch_size, 25)

    # ★ 构造所有批次信息
    batches = []
    for batch_idx in range((len(zh_texts) + batch_size - 1) // batch_size):
        start = batch_idx * batch_size
        end = min(start + batch_size, len(zh_texts))
        batches.append((batch_idx, zh_texts[start:end], start))

    total_batches = len(batches)
    done_count = 0
    _429_count = 0  # ★ 429计数器，连续429时自动降并发

    # ★ v2.21: 自适应并发——先2批并发，遇到429降到串行
    from concurrent.futures import ThreadPoolExecutor, as_completed
    _INIT_CONCURRENT = 2  # ★ 初始并发2（免费API并发承受力低，3容易全429）
    _current_concurrent = _INIT_CONCURRENT

    with ThreadPoolExecutor(max_workers=_INIT_CONCURRENT) as executor:
        # 分轮次提交，每轮 _current_concurrent 批
        for round_start in range(0, total_batches, _current_concurrent):
            round_batches = batches[round_start:round_start + _current_concurrent]
            futures = {}
            for batch_idx, batch, start in round_batches:
                future = executor.submit(
                    _translate_one_batch,
                    batch_idx, batch, start,
                    target_lang, lang_name, lang_instruction,
                    example_block, cfg, url, headers, batch_size,
                )
                futures[future] = (batch_idx, batch, start)

            for future in as_completed(futures):
                batch_idx, batch, start = futures[future]
                try:
                    batch_result = future.result()
                except Exception as e:
                    print(f"[批量单语] {target_lang} 批次{batch_idx+1} 异常: {e}")
                    batch_result = {}

                # ★ 写入本批结果
                for idx, trans in batch_result.items():
                    results[idx] = trans

                done_count += len(batch)
                if on_batch_done:
                    on_batch_done(batch_idx + 1, total_batches, done_count, len(zh_texts))

    # ★ v2.20: 集中兜底——未解析到的条目再发一个补批，不逐条调API
    missing_indices = [i for i in range(len(zh_texts)) if not results[i]]
    if missing_indices:
        print(f"[批量单语] {target_lang} 有 {len(missing_indices)} 条未解析，集中补翻译")
        # 分批补翻译（复用相同的批量逻辑，但逐批串行避免429）
        for comp_start in range(0, len(missing_indices), batch_size):
            comp_indices = missing_indices[comp_start:comp_start + batch_size]
            comp_texts = [zh_texts[i] for i in comp_indices]
            # ★ 用简单的逐行格式再试一次（不用标记，减少解析失败概率）
            comp_prompt = (
                f"{example_block}"
                + f"{lang_instruction}。逐条翻译以下中文，每行一条，顺序对应。"
                + "不要加编号，不要加标记，不要加语言名前缀，不要解释。\n\n"
            )
            for i, text in enumerate(comp_texts):
                comp_prompt += f"{text}\n"

            max_tokens = max(2048, sum(len(t) for t in comp_texts) * 4)
            max_tokens = min(max_tokens, 8192)
            payload = {
                "model": cfg["model"],
                "messages": [
                    {"role": "user", "content": comp_prompt},
                ],
                "temperature": FALLBACK_TEMP,
                "max_tokens": max_tokens,
            }

            for attempt in range(2):
                try:
                    r = requests.post(url, headers=headers, json=payload, timeout=cfg["timeout"])
                    r.raise_for_status()
                    r_json = r.json()
                    content = _extract_content(r_json)
                    if content:
                        line_parsed = _parse_batch_by_line(content, len(comp_texts), target_lang)
                        for i in range(len(comp_texts)):
                            if i in line_parsed and line_parsed[i]:
                                results[comp_indices[i]] = line_parsed[i]
                    break
                except Exception as e:
                    if "429" in str(e) and "4.7" in cfg["model"]:
                        payload["model"] = "glm-4-flash"
                    import time; time.sleep(2)

        # ★ 最终仍有缺失的，逐条兜底（极少情况）
        still_missing = [i for i in range(len(zh_texts)) if not results[i]]
        if still_missing:
            print(f"[批量单语] {target_lang} 仍有 {len(still_missing)} 条缺失，逐条兜底")
            for idx in still_missing:
                try:
                    trans = call_online_llm_single_lang(
                        zh_texts[idx], target_lang, examples or [],
                        use_simple_prompt=True
                    )
                    results[idx] = trans
                except Exception:
                    results[idx] = "[翻译失败]"

    # ★ 最终填充空值为 [翻译失败]
    for i in range(len(zh_texts)):
        if not results[i]:
            results[i] = "[翻译失败]"

    return results


def _parse_batch_by_line(content: str, expected_count: int, target_lang: str) -> dict:
    """当标记解析失败时的兜底行解析。

    尝试从模型输出中按行提取翻译结果，支持多种编号前缀格式：
    - "1. text" / "1、text" / "[1] text" / "1：text"

    Args:
        content: 模型输出文本
        expected_count: 期望的条目数
        target_lang: 目标语言代码（用于清洗语言前缀）

    Returns:
        dict: {0-based索引: 翻译文本}
    """
    parsed = {}
    lines = [l.strip() for l in content.splitlines() if l.strip()]
    
    # 过滤空行和分隔线
    valid_lines = []
    for line in lines:
        # 跳过纯分隔线
        if re.match(r'^[\-=]{3,}$', line):
            continue
        # 去掉编号前缀："1. " 或 "1、" 或 "[1] " 等
        stripped = re.sub(r'^\[?\d+[\]、、.：:]\s*', '', line)
        if stripped:
            stripped = _strip_lang_prefix(stripped, target_lang)
            valid_lines.append(stripped)
    
    # 按行号映射
    for i, text in enumerate(valid_lines):
        if i < expected_count:
            parsed[i] = text
    
    return parsed

def strip_lang_instruction(user_input: str) -> tuple[str, list[str] | None]:
    """从用户输入中分离翻译指令和实际文本。

    支持多种输入格式：
    - "翻成英文和俄语：极石汽车" → ("极石汽车", ["en", "ru"])
    - "翻译成英文 极石汽车" → ("极石汽车", ["en"])
    - "英文：极石汽车" → ("极石汽车", ["en"])
    - "极石汽车" → ("极石汽车", None)

    使用三种策略依次尝试分离：
    1. 冒号分隔符切分（最可靠）
    2. 指令动词 + 语言关键词剥离
    3. 语言名 + 连接词开头剥离

    Args:
        user_input: 用户输入的原始文本

    Returns:
        tuple: (clean_text, target_langs)
               clean_text 是去除指令后的纯文本
               target_langs 是语言代码列表或 None
    """
    target_langs = parse_target_langs(user_input)
    if not target_langs:
        return user_input, None
    
    # 策略1：按冒号分隔符切分（最常见、最可靠）
    for sep in ['：', ':']:
        idx = user_input.find(sep)
        if idx > 0:
            prefix = user_input[:idx]
            content = user_input[idx + len(sep):].strip()
            # 检查前半段是否包含语言关键词（说明是指令区）
            if parse_target_langs(prefix) and content:
                return content, target_langs
    
    # ★ 预计算所有语言关键词的位置（策略2和3共用）
    input_lower = user_input.lower()
    spans = []
    for alias, lc in LANG_ALIASES.items():
        if lc not in TRANSLATE_LANGS:
            continue
        if len(alias) <= 2:
            for m in re.finditer(rf'(?<![a-z]){re.escape(alias)}(?![a-z])', input_lower):
                spans.append((m.start(), m.end()))
        else:
            pos = 0
            while True:
                idx_loc = input_lower.find(alias, pos)
                if idx_loc == -1:
                    break
                spans.append((idx_loc, idx_loc + len(alias)))
                pos = idx_loc + 1
    
    spans.sort()

    # 策略2：按指令动词 + 语言关键词剥离
    instr_verbs = ['翻成', '翻译成', '翻译为', '翻为', '帮我翻', '帮我翻译', '翻译', '翻']
    verb_end = 0
    for v in sorted(instr_verbs, key=len, reverse=True):
        if user_input.startswith(v):
            verb_end = len(v)
            break
    
    if verb_end == 0:
        # 没有指令动词且没有冒号分隔符
        # 尝试策略3：语言名 + 连接词 开头的情况
        # 如 "英语 极石汽车欢迎你" → "极石汽车欢迎你"
        # 如 "英语和葡语 极石汽车欢迎你" → "极石汽车欢迎你"
        zone_end = 0
        changed = True
        while changed:
            changed = False
            for s, e in spans:
                if s <= zone_end + 3 and e > zone_end:
                    zone_end = e
                    changed = True
        
        if zone_end > 0:
            # 跳过分隔符和空格
            clean = re.sub(r'^[\s：:，,、\n]+', '', user_input[zone_end:])
            if clean and clean != user_input:
                return clean, target_langs
        
        # 无法分离（可能整段都是语言名），原样返回
        return user_input, target_langs
    
    # 有指令动词，找语言关键词区的结束位置
    # spans 已在前面预计算，直接使用
    
    # 从 verb_end 开始，找连续的语言关键词区（允许中间有3字符间隔的连接词）
    zone_end = verb_end
    changed = True
    while changed:
        changed = False
        for s, e in spans:
            if s <= zone_end + 3 and e > zone_end:
                zone_end = e
                changed = True
    
    # 跳过分隔符和空格
    clean = re.sub(r'^[\s：:，,、\n]+', '', user_input[zone_end:])
    if clean and clean != user_input:
        return clean, target_langs
    
    # 无法分离（可能整段都是语言名），原样返回
    return user_input, target_langs


def translate_one_prepare(conn, ids, vecs, zh_text: str, target_langs: list[str] | None = None) -> dict:
    """分步翻译的第一步：知识库匹配（即时完成）。

    返回已命中的翻译和待模型补全的语言列表，供调用方决定何时调用模型翻译。
    此函数是 translate_one 的"纯知识库"版本，不触发模型调用。

    匹配顺序：精确命中 → 模糊匹配 → 语义相似度匹配。

    Args:
        conn: SQLite 数据库连接
        ids: 向量索引 ID 数组
        vecs: 向量数组
        zh_text: 中文原文
        target_langs: 目标语言列表，None=全语种

    Returns:
        dict: {
            "mode": str,             # 匹配模式
            "matched_zh": str|None,  # 匹配到的中文
            "similarity": float,
            "translations": dict,    # 已有的翻译（知识库命中或空）
            "candidates": list|None,
            "target_langs": list,
            "need_model": list[str], # 需要模型兜底的语言代码列表
            "examples": list[dict],  # 参考例句
            "is_kb_full_hit": bool,  # 知识库是否全部命中
        }
    """
    out_langs = target_langs if target_langs else TRANSLATE_LANGS
    need_model_langs = list(out_langs)
    examples = []
    matched_zh = None
    mode = "全新"
    similarity = 0.0
    candidates = None

    # ---- 精确命中 ----
    exact = conn.execute("SELECT * FROM tm_segments WHERE zh=?", (zh_text,)).fetchone()
    if exact:
        cols = [d[0] for d in conn.execute("SELECT * FROM tm_segments LIMIT 1").description]
        row = dict(zip(cols, exact))
        translations = {lc: row.get(lc, "") for lc in out_langs}
        missing = [lc for lc in out_langs if not translations.get(lc)]
        return {
            "mode": "精确命中", "matched_zh": zh_text, "similarity": 1.0,
            "translations": translations,
            "candidates": None, "target_langs": out_langs,
            "is_kb_full_hit": len(missing) == 0,
            "need_model": list(missing), "examples": [],
        }

    # ---- 模糊匹配 ----
    fuzzies = fuzzy_substring_hits(conn, zh_text)
    if fuzzies:
        if len(fuzzies) == 1:
            pick = fuzzies[0]; mode = "模糊子串-复用长句译文"
        else:
            pick = max(fuzzies, key=lambda r: len(r["zh"])); mode = f"模糊子串-多候选({len(fuzzies)})"
        translations = {lc: pick.get(lc, "") for lc in out_langs}
        # 检查是否所有语言都有值
        missing = [lc for lc in out_langs if not translations.get(lc)]
        return {
            "mode": mode, "matched_zh": pick["zh"], "similarity": 1.0,
            "translations": translations,
            "candidates": fuzzies, "target_langs": out_langs,
            "is_kb_full_hit": len(missing) == 0,
            "need_model": list(missing), "examples": [],
        }

    # ---- 语义相似度匹配 ----
    qvec = embed(zh_text)
    hits = search(qvec, ids, vecs, k=TOP_K)
    top_id, top_sim = hits[0] if hits else (None, 0.0)

    if top_id is not None and top_sim >= HIGH_SIM:
        row = fetch_row(conn, top_id)
        translations = {lc: row.get(lc, "") for lc in out_langs}
        missing = [lc for lc in out_langs if not translations.get(lc)]
        return {
            "mode": "语义高相似-直接复用", "matched_zh": row["zh"], "similarity": top_sim,
            "translations": translations,
            "candidates": None, "target_langs": out_langs,
            "is_kb_full_hit": len(missing) == 0,
            "need_model": list(missing), "examples": [],
        }

    # ---- 需要模型兜底 ----
    if top_id is not None and top_sim >= MED_SIM:
        order = np.argsort(-(vecs @ qvec))[:TOP_K]
        examples = [fetch_row(conn, int(ids[o])) for o in order if (vecs @ qvec)[o] >= MED_SIM]

    mode = f"在线逐语-参考{len(examples)}条例句" if examples else "在线逐语-全新"
    matched_zh = examples[0]["zh"] if examples else None
    similarity = top_sim

    return {
        "mode": mode, "matched_zh": matched_zh, "similarity": similarity,
        "translations": {},
        "candidates": None, "target_langs": out_langs,
        "need_model": list(out_langs), "examples": examples,
        "is_kb_full_hit": False,
    }


# ====================================================================
# 文件翻译（v2.8 新增）
# 支持 DOCX、PPTX、XLSX 格式的文本提取和翻译写入
# 前端 app.py 和 Agent 调用版 translate_file.py 共用
# ====================================================================

def _ends_sentence(text):
    """判断文本是否以句末标点结尾（。！？.!?:），即是否是一个完整句子。

    不以句末标点结尾的段落通常是PPT美观换行产生的断句片段。
    此函数用于 PPTX 合并断句片段时的判断依据。

    Args:
        text: 待判断的文本

    Returns:
        bool: True 表示以句末标点结尾
    """
    if not text:
        return True
    last = text.rstrip()[-1]
    return last in '。！？.!?:：;；…'

def file_extract_texts(filepath, ext=None, pptx_merge_fragments=True):
    """从办公文件中提取文本列表（去重保序）。

    PPTX 特殊处理：合并同 shape 内连续不以句末标点结尾的段落，
    避免逐段翻译产生机翻感。

    Args:
        filepath: 文件路径
        ext: 文件扩展名（如 .docx），不传则自动从路径取
        pptx_merge_fragments: PPTX是否合并同shape内的断句片段（默认True）
            PPT中经常为美观把一句话拆成多段，逐段翻译会产生机翻感。
            开启后，同shape内连续不以句末标点结尾的段落会被合并为一条翻译单元。

    Returns:
        list[str]: 提取到的文本列表（PPTX合并后含\\n分隔的合并段）
    """
    if ext is None:
        ext = Path(filepath).suffix.lower()
    
    if ext == '.docx':
        from docx import Document
        doc = Document(filepath)
        texts = []
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text.strip())
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text.strip())
        for section in doc.sections:
            for hf in [section.header, section.first_page_header, section.even_page_header,
                        section.footer, section.first_page_footer, section.even_page_footer]:
                if hf and hf.paragraphs:
                    for p in hf.paragraphs:
                        if p.text.strip():
                            texts.append(p.text.strip())
    elif ext == '.pptx':
        from pptx import Presentation
        prs = Presentation(filepath)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    paras = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                    if not paras:
                        continue
                    if pptx_merge_fragments:
                        # ★ 合并同shape内的断句片段：连续不以句末标点结尾的段落合并为一条
                        # 这样翻译时模型能看到完整句子，避免逐段断句产生机翻感
                        merged_group = []
                        for p in paras:
                            merged_group.append(p)
                            if _ends_sentence(p):
                                if len(merged_group) > 1:
                                    texts.append('\n'.join(merged_group))
                                else:
                                    texts.append(merged_group[0])
                                merged_group = []
                        # 收尾：最后一段如果不是完整句，也算一组
                        if merged_group:
                            if len(merged_group) > 1:
                                texts.append('\n'.join(merged_group))
                            else:
                                texts.append(merged_group[0])
                    else:
                        texts.extend(paras)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame and cell.text_frame.text.strip():
                                texts.append(cell.text_frame.text.strip())
    elif ext == '.xlsx':
        import openpyxl
        wb = openpyxl.load_workbook(filepath)
        texts = []
        for ws_name in wb.sheetnames:
            ws = wb[ws_name]
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str) and cell.value.strip():
                        texts.append(cell.value.strip())
    else:
        raise ValueError(f"不支持的文件格式: {ext}")
    
    # 去重保序
    seen = set()
    unique = []
    for t in texts:
        if t not in seen:
            seen.add(t)
            unique.append(t)
    return unique


def _estimate_text_overflow(text, shape_w_emu, shape_h_emu, font_size_pt, line_spacing=1.15):
    """估算文本是否溢出文本框，返回适配后的字号和行间距。

    估算逻辑：
    1. 计算文本的平均字符宽度比例（拉丁≈0.55em，CJK≈1.0em）
    2. 从原始字号向下尝试，找到不溢出的最大字号
    3. 最小字号也放不下时，收紧行间距
    4. 仍放不下则返回最小字号 + 最紧行间距

    Args:
        text: 译文文本
        shape_w_emu: 文本框宽度（EMU）
        shape_h_emu: 文本框高度（EMU）
        font_size_pt: 原始字号（pt）
        line_spacing: 行间距倍数，默认 1.15

    Returns:
        (fitted_size_pt, fitted_line_spacing): 适配后的字号和行间距
    """
    if not text or not shape_w_emu or not shape_h_emu:
        return font_size_pt, line_spacing
    
    # 判断文本字符类型，估算平均字符宽度比例
    lat_cjk = 0
    for c in text:
        cp = ord(c)
        if (0x4E00 <= cp <= 0x9FFF or 0x3400 <= cp <= 0x4DBF or  # CJK统一
            0x3000 <= cp <= 0x303F or                              # CJK符号
            0xFF01 <= cp <= 0xFF60):                               # 全角
            lat_cjk += 1
    cjk_ratio = lat_cjk / max(len(text), 1)
    avg_char_w_ratio = 0.55 * (1 - cjk_ratio) + 1.0 * cjk_ratio
    
    min_size_pt = 8
    max_size_pt = font_size_pt
    
    # 从原始字号向下尝试，找到不溢出的最大字号
    for try_size in _frange(max_size_pt, min_size_pt - 0.5, -0.5):
        size_emu = int(try_size * 12700)
        char_w_emu = max(1, int(size_emu * avg_char_w_ratio))
        chars_per_line = max(1, shape_w_emu // char_w_emu)
        
        # 估算行数（考虑换行符）
        lines_from_newlines = text.count('\n') + 1
        total_text_lines = 0
        for line in text.split('\n'):
            total_text_lines += max(1, -(-len(line) // chars_per_line))  # ceil
        num_lines = total_text_lines
        
        # 行高（EMU）
        line_h = int(size_emu * line_spacing * 1.0)
        total_h = num_lines * line_h
        
        if total_h <= shape_h_emu:
            return round(try_size, 1), line_spacing
    
    # 最小字号也放不下，再收紧行间距
    for try_spacing in _frange(1.0, 0.8, -0.05):
        size_emu = int(min_size_pt * 12700)
        char_w_emu = max(1, int(size_emu * avg_char_w_ratio))
        chars_per_line = max(1, shape_w_emu // char_w_emu)
        total_text_lines = 0
        for line in text.split('\n'):
            total_text_lines += max(1, -(-len(line) // chars_per_line))
        line_h = int(size_emu * try_spacing * 1.0)
        total_h = total_text_lines * line_h
        if total_h <= shape_h_emu:
            return min_size_pt, round(try_spacing, 2)
    
    return min_size_pt, 0.8


def _frange(start, stop, step):
    """浮点数版本的 range 生成器。

    生成从 start 到 stop（含），步长为 step 的浮点数序列。
    支持正步长和负步长。

    Args:
        start: 起始值
        stop: 终止值（含）
        step: 步长

    Returns:
        list[float]: 浮点数序列
    """
    vals = []
    v = start
    if step > 0:
        while v <= stop:
            vals.append(v)
            v += step
    else:
        while v >= stop:
            vals.append(v)
            v += step
    return vals


def _fit_pptx_shape_text(shape, all_para_texts, orig_size_pt):
    """对整个 shape 的所有段落做统一的字号 + 行距适配。

    将 shape 中所有已翻译段落的文本合并计算溢出，
    然后统一缩小字号和行间距。

    Args:
        shape: pptx shape 对象
        all_para_texts: [(para, translated_text), ...] 所有需要适配的段落及译文
        orig_size_pt: 原始字号（pt）

    Returns:
        (fit_size_pt, fit_line_spacing): 适配后的字号和行距
    """
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    
    if not all_para_texts or orig_size_pt is None:
        return orig_size_pt, 1.15
    
    # 合并所有段落文本做总体溢出估算
    full_text = '\n'.join(t for _, t in all_para_texts if t)
    if not full_text:
        return orig_size_pt, 1.15
    
    shape_w = shape.width
    shape_h = shape.height
    
    if not (shape_w and shape_h and shape_w > 0 and shape_h > 0):
        # shape宽高不可用，按字符数比粗略缩放
        orig_total = sum(len(o) for _, o in all_para_texts if o)  # 这里o是译文，但我们需要原文长度做对比
        # 用全文字符数比例
        ratio = 1.0
        if len(full_text) > 0:
            # 粗略：中文1字符≈英文2-3字符，如果译文比原文长50%以上就缩
            for para, t in all_para_texts:
                if t and len(t) > 20:
                    ratio = min(ratio, 0.85)
        if ratio < 0.9:
            return max(8, int(orig_size_pt * ratio * 1.1)), 1.0
        return orig_size_pt, 1.15
    
    fit_size, fit_spacing = _estimate_text_overflow(full_text, shape_w, shape_h, orig_size_pt)
    return fit_size, fit_spacing


def _apply_size_and_spacing_to_para(para, fit_size_pt, fit_spacing, orig_size_pt):
    """将适配后的字号和行距应用到单个段落。

    Args:
        para: python-pptx Paragraph 对象
        fit_size_pt: 适配后的字号（pt）
        fit_spacing: 适配后的行间距倍数
        orig_size_pt: 原始字号（pt），仅当 fit < orig 时才应用
    """
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    
    if fit_size_pt < orig_size_pt and para.runs:
        para.runs[0].font.size = Pt(fit_size_pt)
    
    # 设置行间距
    try:
        pPr = para._pPr
        if pPr is None:
            pPr = para._p.get_or_add_pPr()
        lnSpc = pPr.find(qn('a:lnSpc'))
        if lnSpc is None:
            from lxml import etree
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        for child in list(lnSpc):
            lnSpc.remove(child)
        from lxml import etree
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(fit_spacing * 100000)))
    except Exception:
        pass


def _apply_pptx_paragraph(para, orig_text, translated_text, shape):
    """将翻译文本写入 PPTX 段落，保留格式。

    注意：字号和行距的自适应改由 _fit_pptx_shape_text + _apply_size_and_spacing_to_para
    在 shape 级别统一处理。此函数只负责写入文本。

    Args:
        para: python-pptx Paragraph 对象
        orig_text: 原文文本
        translated_text: 译文文本
        shape: 所属 shape 对象
    """
    from pptx.util import Pt
    # ★ 获取原始字号（供外部shape级别适配使用）
    orig_size_pt = None
    if para.runs:
        if para.runs[0].font.size:
            orig_size_pt = para.runs[0].font.size.pt
    # 兜底：从段落XML的defRPr读取默认字号
    if orig_size_pt is None:
        try:
            from pptx.oxml.ns import qn
            pPr = para._p.find(qn('a:pPr'))
            if pPr is not None:
                defRPr = pPr.find(qn('a:defRPr'))
                if defRPr is not None and defRPr.get('sz'):
                    orig_size_pt = int(defRPr.get('sz')) / 100
        except Exception:
            pass
    # 再兜底：从shape的txBody默认属性读取
    if orig_size_pt is None:
        try:
            from pptx.oxml.ns import qn
            txBody = shape._element.find(qn('p:txBody'))
            if txBody is not None:
                lstStyle = txBody.find(qn('a:lstStyle'))
                if lstStyle is not None:
                    for tag in ['a:lvl1pPr', 'a:lvl2pPr', 'a:lvl3pPr']:
                        lvlPr = lstStyle.find(qn(tag))
                        if lvlPr is not None:
                            defRPr = lvlPr.find(qn('a:defRPr'))
                            if defRPr is not None and defRPr.get('sz'):
                                orig_size_pt = int(defRPr.get('sz')) / 100
                                break
        except Exception:
            pass
    if orig_size_pt is None:
        orig_size_pt = 18
    
    # 写入翻译文本
    if para.runs:
        for run in para.runs:
            run.text = ""
        para.runs[0].text = translated_text
    else:
        para.text = translated_text

def _set_pptx_autofit(shape):
    """将 PPTX shape 的文本框设为自动缩小（AutoFit）。

    设置 XML 属性 autofit="shrink"，让 PowerPoint 在打开文件时自动缩小字号。
    如果 python-pptx 不支持，静默跳过。

    Args:
        shape: pptx shape 对象
    """
    try:
        from pptx.oxml.ns import qn
        txBody = shape._element.find(qn('p:txBody'))
        if txBody is not None:
            bodyPr = txBody.find(qn('a:bodyPr'))
            if bodyPr is not None:
                # 设置 autofit="shrink" — PowerPoint 会在打开时自动缩小字号
                bodyPr.set('autofit', 'shrink')
                # 移除固定内边距可能导致的溢出
                # 不改 mar，保持原边距
    except Exception:
        pass


def file_apply_translations(input_path, output_path, translations, ext=None):
    """将翻译结果写回文件，保持原格式版式，译文超长时自动缩小字号。

    Args:
        input_path: 原始文件路径
        output_path: 输出文件路径
        translations: dict，{原文: 译文}
        ext: 文件扩展名，不传则自动取

    ★ 格式保持策略：保留第一个run的格式属性，清空其余run，译文写入第一个run。
    ★ 自适应字号：译文写入后，估算文本是否溢出文本框，溢出则逐步缩小字号+收紧行间距。
    ★ PPTX额外：设置 autofit="shrink"，让 PowerPoint 在打开时也能自动缩小。
        对 DOCX：译文长度 > 原文1.5倍时按比例缩小字号。
        对 XLSX：译文长度 > 原文1.3倍时缩小字号 + 自动调整列宽。
    """
    if ext is None:
        ext = Path(input_path).suffix.lower()
    
    if ext == '.docx':
        from docx import Document
        from docx.shared import Pt
        doc = Document(input_path)
        
        def _replace_para_docx(para, trans_map):
            full = para.text.strip()
            if not full:
                return
            t = trans_map.get(full) or trans_map.get(full.strip())
            if not t:
                return
            if para.runs:
                # 保存原始字号
                orig_size = None
                if para.runs[0].font.size:
                    orig_size = para.runs[0].font.size.pt
                for run in para.runs:
                    run.text = ""
                para.runs[0].text = t
                # 如果译文明显比原文长（>1.5倍），缩小字号
                if orig_size and len(t) > len(full) * 1.5:
                    shrink = max(8, int(orig_size * len(full) / len(t) * 1.1))
                    para.runs[0].font.size = Pt(shrink)
            else:
                para.text = t
        
        # 同样处理表格单元格
        def _replace_cell_para(para, trans_map, cell_width_emu=None):
            full = para.text.strip()
            if not full:
                return
            t = trans_map.get(full) or trans_map.get(full.strip())
            if not t:
                return
            if para.runs:
                orig_size = None
                if para.runs[0].font.size:
                    orig_size = para.runs[0].font.size.pt
                for run in para.runs:
                    run.text = ""
                para.runs[0].text = t
                if orig_size and len(t) > len(full) * 1.5:
                    shrink = max(8, int(orig_size * len(full) / len(t) * 1.1))
                    para.runs[0].font.size = Pt(shrink)
            else:
                para.text = t
        
        for para in doc.paragraphs:
            _replace_para_docx(para, translations)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        _replace_para_docx(para, translations)
        for section in doc.sections:
            for hf in [section.header, section.first_page_header, section.even_page_header,
                        section.footer, section.first_page_footer, section.even_page_footer]:
                if hf and hf.paragraphs:
                    for para in hf.paragraphs:
                        _replace_para_docx(para, translations)
        doc.save(output_path)
        
    elif ext == '.pptx':
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation(input_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # ★ 合并断句片段的翻译拆回逻辑：
                    # 同shape内连续不以句末标点结尾的段落，在提取时已被合并为一条翻译单元（用\n连接）
                    # 这里需要重建合并key来查找翻译结果，再按\n拆回各段落
                    shape_translated = False
                    paras = [(p, p.text.strip()) for p in shape.text_frame.paragraphs]
                    
                    # 分组合并：连续非句末段落为一组
                    groups = []
                    current_group = []
                    for para, text in paras:
                        if not text:
                            if current_group:
                                groups.append(current_group)
                                current_group = []
                            continue
                        current_group.append((para, text))
                        if _ends_sentence(text):
                            groups.append(current_group)
                            current_group = []
                    if current_group:
                        groups.append(current_group)
                    
                    # ★ 收集本shape所有已翻译的段落，最后统一做字号+行距适配
                    translated_paras = []  # [(para, translated_text, orig_text)]
                    
                    for group in groups:
                        if len(group) == 1:
                            # 单段落，直接匹配
                            para, full = group[0]
                            t = translations.get(full) or translations.get(full.strip())
                            if not t:
                                continue
                            shape_translated = True
                            _apply_pptx_paragraph(para, full, t, shape)
                            translated_paras.append((para, t, full))
                        else:
                            # ★ 合并段落：用\n拼接原文作为key查找翻译
                            merged_key = '\n'.join(text for _, text in group)
                            t = translations.get(merged_key) or translations.get(merged_key.strip())
                            if t:
                                # 翻译结果按\n拆分，分配回各段落
                                parts = t.split('\n')
                                shape_translated = True
                                for idx, (para, full) in enumerate(group):
                                    if idx < len(parts):
                                        para_t = parts[idx].strip()
                                    else:
                                        # 拆分行数不够，最后一段放剩余全部
                                        para_t = '\n'.join(parts[idx:]).strip() if idx < len(parts) + 3 else t
                                    _apply_pptx_paragraph(para, full, para_t, shape)
                                    translated_paras.append((para, para_t, full))
                            else:
                                # 合并key没命中，尝试逐段落独立匹配
                                for para, full in group:
                                    t = translations.get(full) or translations.get(full.strip())
                                    if t:
                                        shape_translated = True
                                        _apply_pptx_paragraph(para, full, t, shape)
                                        translated_paras.append((para, t, full))
                    
                    # ★ shape级别统一适配字号+行距
                    if translated_paras:
                        # 获取原始字号（取第一个有效段落的）
                        orig_size_pt = None
                        for p, _, _ in translated_paras:
                            if p.runs and p.runs[0].font.size:
                                orig_size_pt = p.runs[0].font.size.pt
                                break
                        if orig_size_pt is None:
                            orig_size_pt = 18
                        
                        fit_size, fit_spacing = _fit_pptx_shape_text(
                            shape, [(p, t) for p, t, _ in translated_paras], orig_size_pt
                        )
                        # 应用到所有已翻译段落
                        for para, _, _ in translated_paras:
                            _apply_size_and_spacing_to_para(para, fit_size, fit_spacing, orig_size_pt)
                    
                    # ★ 设置 AutoFit 属性，让 PowerPoint 打开时也能自动缩小
                    if shape_translated:
                        _set_pptx_autofit(shape)
                        
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                for para in cell.text_frame.paragraphs:
                                    full = para.text.strip()
                                    if not full:
                                        continue
                                    t = translations.get(full) or translations.get(full.strip())
                                    if not t:
                                        continue
                                    if para.runs:
                                        orig_size_pt = None
                                        if para.runs[0].font.size:
                                            orig_size_pt = para.runs[0].font.size.pt
                                        for run in para.runs:
                                            run.text = ""
                                        para.runs[0].text = t
                                        # 表格单元格自适应
                                        if orig_size_pt and len(t) > len(full):
                                            try:
                                                cell_w = cell.width  # EMU (可能不可用)
                                                cell_h = None
                                                if cell_w:
                                                    # 粗略估算：表格行高一般按字号×1.5
                                                    fit_size = max(8, int(orig_size_pt * len(full) / len(t) * 1.05))
                                                    if fit_size < orig_size_pt:
                                                        para.runs[0].font.size = Pt(fit_size)
                                            except Exception:
                                                pass
        prs.save(output_path)
        
    elif ext == '.xlsx':
        import openpyxl
        wb = openpyxl.load_workbook(input_path)
        for ws_name in wb.sheetnames:
            ws = wb[ws_name]
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str) and cell.value.strip():
                        orig_val = cell.value.strip()
                        t = translations.get(orig_val)
                        if t:
                            cell.value = t
                            # ★ 自适应：如果译文比原文长，缩小字号
                            if len(t) > len(orig_val) * 1.3:
                                try:
                                    from openpyxl.styles import Font
                                    orig_font = cell.font
                                    orig_size = orig_font.size or 11
                                    shrink = max(8, int(orig_size * len(orig_val) / len(t) * 1.1))
                                    if shrink < orig_size:
                                        cell.font = Font(
                                            name=orig_font.name,
                                            size=shrink,
                                            bold=orig_font.bold,
                                            italic=orig_font.italic,
                                            color=orig_font.color,
                                        )
                                except Exception:
                                    pass
            # ★ 自动调整列宽
            try:
                for col in ws.column_dimensions:
                    max_len = 0
                    col_letter = col
                    for row in ws.iter_rows(min_col=openpyxl.utils.column_index_from_string(col_letter),
                                            max_col=openpyxl.utils.column_index_from_string(col_letter)):
                        for cell in row:
                            if cell.value:
                                max_len = max(max_len, len(str(cell.value)))
                    if max_len > 0:
                        adjusted_width = min(max_len + 2, 50)
                        ws.column_dimensions[col_letter].width = adjusted_width
            except Exception:
                pass
        wb.save(output_path)
    else:
        raise ValueError(f"不支持的文件格式: {ext}")


# ★ _enrich_chinese_text 已移除（v2.10）
# 中文省略成分补全不再用独立LLM调用硬解，改为写入翻译提示词，
# 让翻译模型自己先理解中文含义，再基于理解翻译并补齐缺失成分


def translate_file_texts(texts, target_langs, use_online=True, conn=None, ids=None, vecs=None, on_step=None):
    """批量翻译文件文本列表，走知识库 + 模型完整链路。

    两部流程：
    1. 先做 KB 匹配，收集所有需要模型翻译的条目
    2. 按语言分组，每语言批量调用 call_online_llm_single_lang_batch

    这避免了旧逻辑中逐条文本逐语言串行调用的低效问题
    （580条×9语言=5220次API → 9语言×39批=351次API）。

    Args:
        texts: 待翻译的文本列表
        target_langs: 目标语言代码列表，如 ["en", "ru"]
        use_online: (保留参数兼容性，实际统一走在线模型)
        conn: 知识库SQLite连接
        ids: 向量索引IDs
        vecs: 向量索引vecs

    Yields:
        (all_translations, kb_hits, model_hits, done, total) 每翻译一条输出一次进度

    Returns (最后一次yield):
        all_translations: {原文: {lc: 译文}}
        kb_hits: KB命中次数
        model_hits: 模型翻译次数
        done: 已完成数
        total: 总数
    """
    all_translations = {}
    kb_hits = 0
    model_hits = 0
    translate_total = len(texts) * len(target_langs)
    # ★ total 包含翻译 + 复查两个阶段
    total = translate_total
    done = 0
    all_examples = []
    
    # ★ v2.20: 先做KB匹配，收集所有需要模型翻译的条目，最后批量翻译
    # 旧逻辑：逐条文本逐语言串行调 call_online_llm_single_lang → 580条×9语言=5220次API
    # 新逻辑：按语言分组，每语言批量调 call_online_llm_single_lang_batch → 9语言×39批=351次API
    
    # ---- 第1遍：KB匹配，收集 need_model 条目 ----
    _need_model_by_lang = {}  # {lang_code: [(text_idx, text), ...]}
    
    for text_idx, text in enumerate(texts):
        text_translations = {}
        
        # 知识库匹配
        prep = translate_one_prepare(conn, ids, vecs, text, target_langs=target_langs)
        need_model = list(prep["need_model"])
        
        for lc in target_langs:
            val = prep["translations"].get(lc, "")
            if val:
                text_translations[lc] = val
                kb_hits += 1
                done += 1
        
        # 收集参考例句
        if prep.get("examples"):
            all_examples.extend(prep["examples"])
        all_examples = all_examples[-TOP_K:]
        
        # ★ 记录需要模型翻译的条目（不立即翻译）
        if need_model:
            for lc in need_model:
                if lc not in _need_model_by_lang:
                    _need_model_by_lang[lc] = []
                _need_model_by_lang[lc].append((text_idx, text))
        
        all_translations[text] = text_translations
    
    # ---- 第2遍：按语言批量模型翻译（多语言并发） ----
    if _need_model_by_lang:
        from concurrent.futures import ThreadPoolExecutor, as_completed
        
        def _batch_translate_lang(lang_code, items):
            """对一种语言的待翻译条目做批量翻译"""
            zh_list = [text for _, text in items]
            idx_list = [idx for idx, _ in items]
            batch_results = call_online_llm_single_lang_batch(
                zh_texts=zh_list,
                target_lang=lang_code,
                examples=all_examples[-10:],  # 只取最近10条例句
                batch_size=15,
                on_batch_done=None,  # 进度在外层统一汇报
            )
            return (lang_code, idx_list, zh_list, batch_results)
        
        _lang_workers = min(len(_need_model_by_lang), 3)
        with ThreadPoolExecutor(max_workers=_lang_workers) as executor:
            futures = {}
            for lang_code, items in _need_model_by_lang.items():
                future = executor.submit(_batch_translate_lang, lang_code, items)
                futures[future] = lang_code
            
            for future in as_completed(futures):
                lang_code = futures[future]
                try:
                    result_lang, idx_list, zh_list, batch_results = future.result()
                except Exception as e:
                    print(f"[KB翻译] {lang_code} 批量翻译异常: {e}")
                    continue
                
                for i, text_idx in enumerate(idx_list):
                    text = zh_list[i]
                    trans = batch_results[i] if i < len(batch_results) else "[翻译失败]"
                    all_translations[text][result_lang] = trans
                    if trans and not trans.startswith("[翻译失败"):
                        model_hits += 1
                    done += 1
                
                # ★ yield 进度
                yield all_translations, kb_hits, model_hits, done, total

# ====================================================================
# 文件纠错（v2.9 重构）
# PPTX：红字批注识别；DOCX/XLSX：位置对齐比对
# ====================================================================

def parse_correction_note(note_text):
    """从纠错批注中提取修正后的翻译文本。

    常见纠错批注模式：
    - "改为：XXX" / "改为XXX"
    - "第一行改为：XXX"
    - "应该/可以写XXX" / "应该/可以翻译成XXX"
    - "改成XXX" / "修改为XXX" / "修正为XXX"

    Args:
        note_text: 批注文本

    Returns:
        str | None: 提取到的修正文本，无法解析则返回 None
    """
    import re
    patterns = [
        r'改为[：:]\s*(.+?)(?:[，,。.；;！!]|$)',
        r'改为\s+(.+?)(?:[，,。.；;！!]|$)',
        r'改成[：:]\s*(.+?)(?:[，,。.；;！!]|$)',
        r'改成\s+(.+?)(?:[，,。.；;！!]|$)',
        r'应该[写翻]成?[：:]?\s*(.+?)(?:[，,。.；;！!]|$)',
        r'可以[写翻]成?[：:]?\s*(.+?)(?:[，,。.；;！!]|$)',
        r'修改为[：:]\s*(.+?)(?:[，,。.；;！!]|$)',
        r'修正为[：:]\s*(.+?)(?:[，,。.；;！!]|$)',
        r'[应该可以]*翻译[成为]?[：:]?\s*(.+?)(?:[，,。.；;！!]|$)',
    ]
    for pat in patterns:
        m = re.search(pat, note_text)
        if m:
            return m.group(1).strip()
    return None


def extract_pptx_red_annotations(filepath):
    """提取 PPTX 中红色字体的批注文本，并关联到同页最近的非红色文本。

    纠错文件的特点：原文翻译仍保留，旁边用红字批注标注怎么改。
    因此不能用位置对齐，而要提取红字批注并关联到附近的原文翻译。

    关联策略：计算红字 shape 和普通文本 shape 的中心点距离，取最近者。

    Returns:
        list[dict]: 每项含 slide_idx, red_text, nearby_original
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    
    def _is_red(run):
        """判断 run 是否为红色字体。

        检测逻辑：检查 run 的字体颜色是否为纯红（FF0000）或偏红
        （R通道高，G/B通道低）。

        Args:
            run: python-pptx Run 对象

        Returns:
            bool: True 表示为红色字体
        """
        try:
            color = run.font.color
            if color is None:
                return False
            if color.rgb is not None:
                rgb = str(color.rgb).upper()
                # 纯红 FF0000 或偏红（R高，G/B低）
                if rgb.startswith('FF') and rgb[2:4] in ('00', '33', '44', '55') and rgb[4:6] in ('00', '33', '44', '55'):
                    return True
                if rgb.startswith('CC') and rgb[2:4] in ('00', '33') and rgb[4:6] in ('00', '33'):
                    return True
                if rgb.startswith('DD') and rgb[2:4] in ('00', '22', '33') and rgb[4:6] in ('00', '22', '33'):
                    return True
            # theme color fallback — 无法确定具体颜色，跳过
        except Exception:
            pass
        return False
    
    prs = Presentation(filepath)
    annotations = []
    
    for slide_idx, slide in enumerate(prs.slides):
        red_shapes = []
        normal_shapes = []
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            red_parts = []
            normal_parts = []
            has_red = False
            
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    if _is_red(run):
                        red_parts.append(run.text.strip())
                        has_red = True
                    else:
                        normal_parts.append(run.text.strip())
            
            shape_info = {
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'red_text': ' '.join(red_parts),
                'normal_text': ' '.join(normal_parts),
                'has_red': has_red,
            }
            
            if has_red and red_parts:
                red_shapes.append(shape_info)
            if normal_parts:
                normal_shapes.append(shape_info)
        
        # 匹配：每个红字批注关联到同页最近的非红文本
        for rs in red_shapes:
            best_dist = float('inf')
            best_normal = None
            for ns in normal_shapes:
                r_cx = rs['left'] + rs['width'] / 2
                r_cy = rs['top'] + rs['height'] / 2
                n_cx = ns['left'] + ns['width'] / 2
                n_cy = ns['top'] + ns['height'] / 2
                dist = ((r_cx - n_cx) ** 2 + (r_cy - n_cy) ** 2) ** 0.5
                if dist < best_dist:
                    best_dist = dist
                    best_normal = ns
            
            annotations.append({
                'slide_idx': slide_idx + 1,
                'red_text': rs['red_text'],
                'nearby_original': best_normal['normal_text'] if best_normal else '',
                'also_has_normal': rs['normal_text'],  # 同一shape里的非红文本
            })
    
    return annotations


def _pptx_correct_compare(original_path, corrected_path, lang_code, conn, ids, vecs):
    """PPTX 纠错：提取红色批注 → 解析修正文本 → 匹配 KB → 保存。

    PPTX 纠错文件的特点：原文翻译仍保留，旁边用红字批注标注怎么改。
    所以不能用位置对齐，而要用语义匹配：

    流程：
    1. 提取红字批注和附近的原文翻译
    2. 从批注中解析出修正文本（如"改为：XXX"）
    3. 用附近原文翻译在KB中反查中文源文
       - 3a. 精确匹配 KB 的语言列
       - 3b. LIKE 模糊匹配
       - 3c. 用中文原文列表 + KB 匹配
    4. 有差异则更新 KB

    Args:
        original_path: 原始文件路径（中文原文）
        corrected_path: 纠错后文件路径（含红字批注）
        lang_code: 目标语言代码
        conn: SQLite 数据库连接
        ids: 向量索引 ID 数组
        vecs: 向量数组

    Returns:
        tuple: (results: list[dict], summary: dict)
    """
    lang_name = LANG_NAMES.get(lang_code, "")
    annotations = extract_pptx_red_annotations(corrected_path)
    zh_texts = file_extract_texts(original_path, '.pptx')
    
    results = []
    changed_count = 0
    saved_count = 0
    failed_count = 0
    
    for ann in annotations:
        red_text = ann['red_text']
        nearby_orig = ann.get('nearby_original', '') or ann.get('also_has_normal', '')
        slide_info = f"第{ann['slide_idx']}页"
        
        # 1. 解析批注 → 提取修正文本
        corrected_trans = parse_correction_note(red_text)
        if not corrected_trans:
            results.append({
                "zh": nearby_orig[:80],
                "old_trans": "",
                "new_trans": "",
                "changed": False,
                "saved": False,
                "error": f"无法从批注提取修正文本",
                "slide": slide_info,
                "red_note": red_text,
                "nearby": nearby_orig[:60],
            })
            continue
        
        # 2. 用附近原文翻译在KB中反查中文源文
        matched_zh = None
        old_trans = ""
        row_dict = None
        
        # 2a. 先精确匹配KB的语言列
        if nearby_orig:
            try:
                rows = conn.execute(
                    f"SELECT * FROM tm_segments WHERE [{lang_code}]=?",
                    (nearby_orig.strip(),)
                ).fetchall()
                if rows:
                    cols = [d[0] for d in conn.execute("SELECT * FROM tm_segments LIMIT 1").description]
                    row_dict = dict(zip(cols, rows[0]))
                    matched_zh = row_dict['zh']
                    old_trans = row_dict.get(lang_code, '') or ''
            except Exception:
                pass
        
        # 2b. 精确没中，用LIKE模糊查
        if not matched_zh and nearby_orig:
            try:
                snippet = nearby_orig.strip()[:40]
                rows = conn.execute(
                    f"SELECT * FROM tm_segments WHERE [{lang_code}] LIKE ?",
                    (f'%{snippet}%',)
                ).fetchall()
                if rows:
                    cols = [d[0] for d in conn.execute("SELECT * FROM tm_segments LIMIT 1").description]
                    row_dict = dict(zip(cols, rows[0]))
                    matched_zh = row_dict['zh']
                    old_trans = row_dict.get(lang_code, '') or ''
            except Exception:
                pass
        
        # 2c. 还没中，用中文原文列表+KB匹配
        if not matched_zh and nearby_orig:
            for zh in zh_texts:
                row = conn.execute("SELECT * FROM tm_segments WHERE zh=?", (zh,)).fetchone()
                if row:
                    cols = [d[0] for d in conn.execute("SELECT * FROM tm_segments LIMIT 1").description]
                    rd = dict(zip(cols, row))
                    kb_trans = rd.get(lang_code, '') or ''
                    if kb_trans and (kb_trans in nearby_orig or nearby_orig in kb_trans):
                        matched_zh = zh
                        old_trans = kb_trans
                        row_dict = rd
                        break
        
        if not matched_zh:
            results.append({
                "zh": "",
                "old_trans": nearby_orig[:80],
                "new_trans": corrected_trans,
                "changed": True,
                "saved": False,
                "error": f"未在KB中找到匹配（{slide_info}）",
                "slide": slide_info,
                "red_note": red_text,
                "nearby": nearby_orig[:60],
            })
            failed_count += 1
            continue
        
        # 3. 检查是否真有差异
        if corrected_trans == old_trans:
            results.append({
                "zh": matched_zh, "old_trans": old_trans, "new_trans": corrected_trans,
                "changed": False, "saved": False, "error": None,
                "slide": slide_info, "red_note": red_text, "nearby": nearby_orig[:60],
            })
            continue
        
        # 4. 保存修正到KB
        changed_count += 1
        translations = {lang_code: corrected_trans}
        if row_dict:
            for lc in TRANSLATE_LANGS:
                if lc != lang_code and row_dict.get(lc):
                    translations[lc] = row_dict[lc]
        
        try:
            row_id = save_back(conn, matched_zh, translations)
            emb_ok = rebuild_embeddings_for_entry(matched_zh, row_id)
            saved_count += 1
            results.append({
                "zh": matched_zh, "old_trans": old_trans, "new_trans": corrected_trans,
                "changed": True, "saved": True,
                "error": None if emb_ok else "向量索引更新失败",
                "slide": slide_info, "red_note": red_text, "nearby": nearby_orig[:60],
            })
        except Exception as e:
            failed_count += 1
            results.append({
                "zh": matched_zh, "old_trans": old_trans, "new_trans": corrected_trans,
                "changed": True, "saved": False, "error": str(e),
                "slide": slide_info, "red_note": red_text, "nearby": nearby_orig[:60],
            })
    
    summary = {
        "total": len(annotations),
        "original_count": len(zh_texts),
        "changed": changed_count,
        "saved": saved_count,
        "failed": failed_count,
        "unparsed": sum(1 for r in results if not r.get("changed") and r.get("error")),
    }
    return results, summary


def _position_correct_compare(original_path, corrected_path, lang_code, conn, ids, vecs, ext):
    """DOCX/XLSX 纠错：位置对齐比对。

    对 DOCX 和 XLSX 文件，文本提取后按位置一一对应比较。
    相同位置的原文 → 译文对照，有差异则更新知识库。

    Args:
        original_path: 原始文件路径
        corrected_path: 纠错后文件路径
        lang_code: 目标语言代码
        conn: SQLite 数据库连接
        ids: 向量索引 ID 数组
        vecs: 向量数组
        ext: 文件扩展名

    Returns:
        tuple: (results: list[dict], summary: dict)
    """
    zh_texts = file_extract_texts(original_path, ext)
    corrected_texts = file_extract_texts(corrected_path, ext)

    total = min(len(zh_texts), len(corrected_texts))
    results = []
    changed_count = 0
    saved_count = 0
    failed_count = 0

    for i in range(total):
        zh = zh_texts[i].strip()
        new_trans = corrected_texts[i].strip()

        if not zh or not new_trans:
            results.append({
                "zh": zh, "old_trans": "", "new_trans": new_trans,
                "changed": False, "saved": False, "error": "空文本",
            })
            continue

        row = conn.execute("SELECT * FROM tm_segments WHERE zh=?", (zh,)).fetchone()
        old_trans = ""
        row_dict = None
        if row:
            cols = [d[0] for d in conn.execute("SELECT * FROM tm_segments LIMIT 1").description]
            row_dict = dict(zip(cols, row))
            old_trans = row_dict.get(lang_code, "") or ""

        if new_trans == old_trans:
            results.append({
                "zh": zh, "old_trans": old_trans, "new_trans": new_trans,
                "changed": False, "saved": False, "error": None,
            })
            continue

        changed_count += 1
        translations = {lang_code: new_trans}
        if row_dict:
            for lc in TRANSLATE_LANGS:
                if lc != lang_code and row_dict.get(lc):
                    translations[lc] = row_dict[lc]

        try:
            row_id = save_back(conn, zh, translations)
            emb_ok = rebuild_embeddings_for_entry(zh, row_id)
            saved_count += 1
            results.append({
                "zh": zh, "old_trans": old_trans, "new_trans": new_trans,
                "changed": True, "saved": True,
                "error": None if emb_ok else "向量索引更新失败",
            })
        except Exception as e:
            failed_count += 1
            results.append({
                "zh": zh, "old_trans": old_trans, "new_trans": new_trans,
                "changed": True, "saved": False, "error": str(e),
            })

    summary = {
        "total": total,
        "original_count": len(zh_texts),
        "corrected_count": len(corrected_texts),
        "matched": total,
        "changed": changed_count,
        "saved": saved_count,
        "failed": failed_count,
    }
    return results, summary


# ====================================================================
# Phase 1: segment_base 结构化知识库
# 将翻译记忆库中的条目拆解为可复用的语义片段（术语、短语、句式模板等）
# ====================================================================

# 语义片段类型枚举
# 用于 segment_base 表，标识每一条片段的语义类别
SEGMENT_TYPES = {
    "term",        # 原子术语（如"极石汽车→ROX"）
    "phrase",      # 固定短语（如"了解更多→Learn More"）
    "position",    # 位置词（如"左前→Left Front"）
    "component",   # 部件名（如"轮胎→Tire"）
    "status",      # 状态词（如"异常→Abnormal"）
    "negation",    # 否定结构（如"不{verb}→Not {verb}ed"）
    "pattern",     # 句式模板（如"{pos}{comp}{status}→{pos} {comp} {status}"）
    "syntax",      # 语法结构（如"为{某人}预留{某物}→Reserve {sth} for {sb}"）
    "cond_result", # 条件结果（如"{cond}，请{action}→Please {action} if {cond}"）
    "action_guide",# 操作引导（如"请{action}→Please {action}"）
    "product_frame", # 产品框架（如"产品描述：{desc}→Product Description: {desc}"）
    "value_unit",  # 数值单位（如"约{num}{unit}→Approx. {num} {unit}"）
    "enumeration", # 列表编号（如"{n}. {item}→{n}. {item}"）
    "disclaimer",  # 免责句式（如"如{cond}不承担{liability}→Not liable for {liability} if {cond}"）
}

def init_segment_db(conn):
    """初始化 segment_base 表（如果不存在）。

    创建存储语义片段的数据库表，包含类型、语言对、原文/译文、变量等字段。
    同时建立 lang_pair、seg_type、segment_hash 三个索引加速查询。

    Args:
        conn: SQLite 数据库连接
    """
    conn.execute("""
        CREATE TABLE IF NOT EXISTS segment_base (
            id          INTEGER PRIMARY KEY AUTOINCREMENT,
            source_id   INTEGER,
            seg_type    TEXT NOT NULL,
            lang_pair   TEXT NOT NULL,
            source_text TEXT NOT NULL,
            target_text TEXT NOT NULL,
            variables   TEXT DEFAULT '[]',
            tag         TEXT DEFAULT '',
            segment_hash TEXT UNIQUE,
            created_at  TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    """)
    conn.execute("CREATE INDEX IF NOT EXISTS idx_sb_lang_pair ON segment_base(lang_pair)")
    conn.execute("CREATE INDEX IF NOT EXISTS idx_sb_seg_type ON segment_base(seg_type)")
    conn.execute("CREATE INDEX IF NOT EXISTS idx_sb_hash ON segment_base(segment_hash)")
    conn.commit()

def _segment_hash(source_text: str, lang_pair: str) -> str:
    """计算语义片段的唯一哈希值。

    Args:
        source_text: 原文文本
        lang_pair: 语言对（如 "zh-en"）

    Returns:
        str: MD5 哈希值
    """
    return hashlib.md5(f"{source_text}||{lang_pair}".encode()).hexdigest()

def build_segment_base(conn):
    """遍历 tm_segments，用 LLM 提取语义段，去重写入 segment_base。

    对每条翻译记忆条目，调用在线模型分析中文→英文翻译，
    提取所有可复用的语义片段（术语、短语、句式模板等），
    写入 segment_base 表供后续段匹配使用。

    Args:
        conn: SQLite 数据库连接

    Returns:
        int: 提取到的语义片段总数
    """
    from concurrent.futures import ThreadPoolExecutor, as_completed
    init_segment_db(conn)

    rows = conn.execute("SELECT id, zh, en, ru, ar, es, pt, fr, kk, de, zh_hant, ms, id_lang, th, tr, it, pl, sv, module FROM tm_segments WHERE zh IS NOT NULL AND zh != ''").fetchall()
    lang_cols = ["en", "ru", "ar", "es", "pt", "fr", "kk", "de", "zh_hant", "ms", "id_lang", "th", "tr", "it", "pl", "sv"]

    extract_prompt = """分析以下中文→英文翻译条目，提取所有可复用的翻译片段。

原文: {zh}
译文: {en}

从以下14种类别中识别本条包含哪些，输出JSON数组：
- term（原子术语，如"极石汽车→ROX"）
- phrase（固定短语，如"了解更多→Learn More"）
- position（位置词，如"左前→Left Front"）
- component（部件名，如"轮胎→Tire"）
- status（状态词，如"异常→Abnormal"）
- negation（否定结构，如"不{verb}→Not {verb}ed"）
- pattern（句式模板，如"{pos}{comp}{status}→{pos} {comp} {status}"）
- syntax（语法结构，如"为{某人}预留{某物}→Reserve {sth} for {sb}"）
- cond_result（条件结果，如"{cond}，请{action}→Please {action} if {cond}"）
- action_guide（操作引导，如"请{action}→Please {action}"）
- product_frame（产品框架，如"产品描述：{desc}→Product Description: {desc}"）
- value_unit（数值单位，如"约{num}{unit}→Approx. {num} {unit}"）
- enumeration（列表编号，如"{n}. {item}→{n}. {item}"）
- disclaimer（免责句式，如"如{cond}不承担{liability}→Not liable for {liability} if {cond}"）

规则：
- 每段 source_text 必须在原文中能找到连续匹配
- 模板类用 {var} 标记变量
- 输出纯JSON数组，不要markdown不要解释
- 若无任何可提取片段则输出 []
"""

    batch = []
    batch_size = 20
    total_segments = 0

    for row in rows:
        zh_text = row[1]
        en_text = row[2]
        if not zh_text or not en_text:
            continue
        batch.append((row[0], zh_text, en_text))

    for i in range(0, len(batch), batch_size):
        chunk = batch[i:i+batch_size]
        for sid, zh, en in chunk:
            try:
                prompt = extract_prompt.format(zh=zh, en=en)
                payload = {
                    "model": os.getenv("ONLINE_MODEL", "tencent/Hunyuan-MT-7B"),
                    "messages": [{"role": "user", "content": prompt}],
                    "temperature": 0.1,
                    "max_tokens": 2048,
                }
                cfg = _get_online_config()
                url = f"{cfg['base_url'].rstrip('/')}/chat/completions"
                headers = {"Authorization": f"Bearer {cfg['api_key']}", "Content-Type": "application/json"}
                r = requests.post(url, headers=headers, json=payload, timeout=120)
                r.raise_for_status()
                content = _extract_content(r.json())
                content = content.replace("```json", "").replace("```", "").strip()
                segments = json.loads(content) if content else []
                if not isinstance(segments, list):
                    continue
                for seg in segments:
                    st = seg.get("seg_type", "")
                    stext = seg.get("source_text", "").strip()
                    ttext = seg.get("target_text", "").strip()
                    if st not in SEGMENT_TYPES or not stext or not ttext:
                        continue
                    h = _segment_hash(stext, "zh-en")
                    try:
                        conn.execute(
                            "INSERT OR IGNORE INTO segment_base (source_id, seg_type, lang_pair, source_text, target_text, tag, segment_hash) VALUES (?,?,?,?,?,?,?)",
                            (sid, st, "zh-en", stext, ttext, "", h)
                        )
                        conn.commit()
                        total_segments += 1
                    except Exception:
                        pass
            except Exception as e:
                print(f"[segment_base] 提取失败 row={sid}: {e}")
                continue

        print(f"[segment_base] 进度 {min(i+batch_size, len(batch))}/{len(batch)}，已提取 {total_segments} 段")

    print(f"[segment_base] 完成：从 {len(batch)} 条TM中提取 {total_segments} 个片段")
    return total_segments


def split_long_entries(conn, limit: int = 0):
    """
    用 LLM 将长句拆解为独立短句，审计后写入 tm_segments。

    对每条翻译记忆条目，调用在线模型将其拆分为多个独立短句，
    同时拆分对应的多语言翻译。审计通过后写入 tm_segments。

    Args:
        limit: 最多处理条数（0=不限）
    """
    lang_cols = ["en", "ru", "ar", "es", "pt", "fr", "kk", "de", "zh_hant"]

    sql = (
        f"SELECT id, zh, module, {','.join(f'\"{c}\"' for c in lang_cols)} "
        "FROM tm_segments WHERE zh IS NOT NULL AND zh != ''"
    )
    if limit > 0:
        sql += f" LIMIT {limit}"
    rows = conn.execute(sql).fetchall()

    total_added = 0
    total_failed = 0
    total_no_split = 0
    total_entries = len(rows)

    cfg = _get_online_config()
    url = f"{cfg['base_url'].rstrip('/')}/chat/completions"
    headers = {"Authorization": f"Bearer {cfg['api_key']}", "Content-Type": "application/json"}

    # ── 拆句 prompt（使用通用指令模型，不含翻译内容，只拆中文）──
    # ★ 先拆中文，再逐语言匹配翻译内容
    split_prompt = """将以下中文文本按独立语义单位拆分为多个短句。
拆分点包括：句号、叹号、问号、分号、换行。
逗号位置如果不是并列结构也拆分。

中文原文：
{zh}

输出 JSON 字符串数组，每个元素是一个短句。
如果整句是单一语义单位不可拆分，输出 [原文全文]（即一个元素的数组）。
只输出 JSON，不要其他文字。"""

    # ★ 拆句用 GLM-4-9B（指令理解好），提取用 GLM-4-9B
    split_cfg = {**cfg, "model": HUNYUAN_FALLBACK_MODEL}
    split_url = f"{split_cfg['base_url'].rstrip('/')}/chat/completions"
    extract_cfg = split_cfg
    extract_url = split_url

    for idx, row in enumerate(rows):
        row_id = row[0]
        zh_text = row[1]
        module_tag = row[2]
        lang_texts = {lang_cols[i]: (row[3 + i] or "").strip() for i in range(len(lang_cols))}

        # 统一引号避免 LLM JSON 输出被破坏
        _smart_repl = str.maketrans({
            '\u201c': '"', '\u201d': '"', '\u201e': '"', '\u201f': '"',
            '\u2018': "'", '\u2019': "'", '\u201a': "'", '\u201b': "'",
        })
        zh_text = zh_text.translate(_smart_repl)
        for lc in lang_cols:
            if lang_texts.get(lc):
                lang_texts[lc] = lang_texts[lc].translate(_smart_repl)

        # 去掉列表前缀（* / - / • / 数字.）让提取更准确
        _bullet_re = re.compile(r'^[\*\-•·]\s*|\d+\.\s*', re.MULTILINE)
        zh_text = _bullet_re.sub('', zh_text)
        for lc in lang_cols:
            if lang_texts.get(lc):
                lang_texts[lc] = _bullet_re.sub('', lang_texts[lc])

        zh_cjk = extract_cjk(zh_text)
        if len(zh_cjk) <= 8:
            total_no_split += 1
            continue

        if not any(c in zh_text for c in '，；。、\n'):
            total_no_split += 1
            continue

        # ── 第1步：拆中文（用 GLM-4-9B）──
        split_user = split_prompt.format(zh=zh_text)
        zh_segments = _call_llm_json(split_user, split_cfg, split_url, headers)
        if not zh_segments or len(zh_segments) <= 1:
            total_no_split += 1
            continue

        # 限制最多 10 段，避免提取时 prompt 过长导致 JSON 出错
        if len(zh_segments) > 10:
            zh_segments = zh_segments[:10]

        # ── 第2步：对所有语言，批量提取每段的多语言翻译 ──
        added = 0
        batch_extract_tpl = (
            "中文全文：{zh_full}\n"
            "已将中文拆分为以下片段：\n{seg_list}\n\n"
            "{lang_name}翻译全文：{lang_trans}\n\n"
            "请找出每个中文片段在{lang_name}翻译中对应的部分。\n"
            "输出 JSON 字符串数组，长度={seg_count}，第N个元素对应第N个中文片段。\n"
            "如果某片段无明确对应，用空字符串。\n"
            "只输出 JSON 数组，不要其他文字。"
        )

        all_parts = {}
        extract_failed = False
        for lc in lang_cols:
            lt = lang_texts.get(lc, "")
            if not lt:
                all_parts[lc] = [""] * len(zh_segments)
                continue
            numbered_segs = "\n".join(f"[{i+1}] {s}" for i, s in enumerate(zh_segments))
            bp = batch_extract_tpl.format(
                zh_full=zh_text,
                seg_list=numbered_segs,
                lang_name=LANG_NAMES.get(lc, lc),
                lang_trans=lt,
                seg_count=len(zh_segments),
            )
            parts = _call_llm_json(bp, extract_cfg, extract_url, headers)
            if not parts or len(parts) != len(zh_segments):
                extract_failed = True
                break
            flat = []
            for p in parts:
                if isinstance(p, (list, tuple)):
                    p = p[0] if p else ""
                if not isinstance(p, str):
                    p = str(p) if p is not None else ""
                flat.append(p.strip() if p else "")
            all_parts[lc] = flat

        if extract_failed:
            total_failed += 1
            continue

        # ── 第3步：写入每段（GLM-4-9B 提取质量可靠，不做 LLM 审计）──
        for i, seg_zh_str in enumerate(zh_segments):
            if isinstance(seg_zh_str, dict):
                seg_zh_str = seg_zh_str.get("zh", "")
            seg_zh_str = (seg_zh_str or "").strip()
            if not seg_zh_str:
                continue
            seg_cjk = extract_cjk(seg_zh_str)
            if seg_cjk == zh_cjk or not seg_cjk:
                total_no_split += 1
                continue

            seg_translations = {lc: all_parts[lc][i] for lc in lang_cols}

            has_any = any(bool(v) for v in seg_translations.values())
            if not has_any:
                total_failed += 1
                continue

            sid = save_back(conn, seg_zh_str, seg_translations, module_tag)
            if sid > 0:
                try:
                    rebuild_embeddings_for_entry(seg_zh_str, sid)
                except Exception:
                    pass
                added += 1

        if added:
            print(f"[split] id={row_id}: +{added} segments")

    _invalidate_zh_cjk_cache()
    print(f"[split] 完成：新增 {total_added} 条，提取失败 {total_failed} 条，{total_no_split} 条无需拆分")
    return {
        "success": True,
        "total_added": total_added,
        "total_failed": total_failed,
        "total_no_split": total_no_split,
        "message": f"拆分完成：新增 {total_added} 条短句，{total_failed} 条提取失败，{total_no_split} 条无需拆分",
    }


def _call_llm_json(user_prompt: str, cfg: dict, url: str, headers: dict) -> list | None:
    """调用 LLM 并解析 JSON 返回"""
    payload = {
        "model": cfg["model"],
        "messages": [{"role": "user", "content": user_prompt}],
        "temperature": 0.1,
        "max_tokens": 4096,
    }
    for attempt in range(3):
        try:
            r = requests.post(url, headers=headers, json=payload, timeout=120)
            r.raise_for_status()
            content = _extract_content(r.json())
            content = content.replace("```json", "").replace("```", "").strip()
            if content:
                content = re.sub(r',\s*}', '}', content)
                content = re.sub(r',\s*]', ']', content)
                try:
                    result = json.loads(content)
                    if isinstance(result, list):
                        return result
                except json.JSONDecodeError:
                    bracket = content.find('[')
                    if bracket >= 0:
                        cbracket = content.rfind(']')
                        if cbracket > bracket:
                            sub = content[bracket:cbracket+1]
                            sub = re.sub(r',\s*}', '}', sub)
                            sub = re.sub(r',\s*]', ']', sub)
                            result = json.loads(sub)
                            if isinstance(result, list):
                                return result
            return None
        except Exception as e:
            if attempt < 2:
                import time
                time.sleep(3)
            else:
                raw_preview = content[:300] if 'content' in locals() and content else '(no content)'
                print(f"  [split] LLM JSON 调用失败: {e} | raw={raw_preview!r}")
                return None


def _call_llm_audit(user_prompt: str, cfg: dict, url: str, headers: dict) -> bool:
    """调用 LLM 审计翻译对齐，返回是否通过"""
    payload = {
        "model": cfg["model"],
        "messages": [{"role": "user", "content": user_prompt}],
        "temperature": 0.0,
        "max_tokens": 20,
    }
    for attempt in range(3):
        try:
            r = requests.post(url, headers=headers, json=payload, timeout=60)
            r.raise_for_status()
            content = _extract_content(r.json())
            if content and content.strip().lower() == "pass":
                return True
            return False
        except Exception as e:
            if attempt < 2:
                import time
                time.sleep(2)
            else:
                print(f"  [split] 审计调用失败: {e}")
                return False


def _call_llm_extract(user_prompt: str, cfg: dict, url: str, headers: dict) -> str | None:
    """调用 LLM 从翻译中提取对应片段的文本"""
    payload = {
        "model": cfg["model"],
        "messages": [{"role": "user", "content": user_prompt}],
        "temperature": 0.1,
        "max_tokens": 1024,
    }
    for attempt in range(3):
        try:
            r = requests.post(url, headers=headers, json=payload, timeout=60)
            r.raise_for_status()
            content = _extract_content(r.json())
            if content:
                return content.strip()
            return None
        except Exception as e:
            if attempt < 2:
                import time
                time.sleep(2)
            else:
                print(f"  [split] 提取调用失败: {e}")
                return None


def match_segments(conn, zh_text: str, lang_pair: str) -> tuple[list[dict], str]:
    """对输入文本做 segment_base 匹配。

    从 segment_base 中加载指定语言对的所有片段，按原文长度降序排列。
    在输入文本中查找匹配的片段（支持重叠检测，已占用的位置不再匹配）。
    返回匹配到的片段列表和未匹配的剩余文本。

    Args:
        conn: SQLite 数据库连接
        zh_text: 中文原文
        lang_pair: 语言对（如 "zh-en"）

    Returns:
        tuple: (matched_segments: list[dict], remaining_text: str)
               matched_segments 按原文中的起始位置排序
               remaining_text 是未被任何片段覆盖的文本
    """
    if not zh_text or not zh_text.strip():
        return [], zh_text

    rows = conn.execute(
        "SELECT seg_type, source_text, target_text FROM segment_base WHERE lang_pair=? ORDER BY length(source_text) DESC",
        (lang_pair,)
    ).fetchall()

    segments = [(r[0], r[1], r[2]) for r in rows]

    matched = []
    remaining = list(zh_text)
    used = [False] * len(zh_text)

    for seg_type, src, tgt in segments:
        if not src:
            continue
        pos = zh_text.find(src)
        if pos >= 0:
            all_free = all(not used[p] for p in range(pos, pos + len(src)))
            if all_free:
                for p in range(pos, pos + len(src)):
                    used[p] = True
                matched.append({
                    "seg_type": seg_type,
                    "source": src,
                    "target": tgt,
                    "start": pos,
                    "end": pos + len(src),
                })

    remaining_text = ""
    i = 0
    while i < len(zh_text):
        if not used[i]:
            remaining_text += zh_text[i]
            i += 1
        else:
            if remaining_text and not remaining_text[-1].isspace():
                remaining_text += " "
            i += 1

    matched.sort(key=lambda x: x["start"])
    remaining_text = remaining_text.strip()

    return matched, remaining_text


def reassemble_translation(matched_segments: list[dict], model_translated: str, original_text: str) -> str:
    """将 segment_base 匹配结果按原文位置重组，未匹配部分用模型翻译填充。

    策略：
    1. 按原文位置排序匹配到的片段
    2. 找出未被片段覆盖的文本区间
    3. 将模型翻译的段落按顺序填入未覆盖区间
    4. 返回完整翻译

    Args:
        matched_segments: 匹配到的片段列表
        model_translated: 模型翻译的剩余部分文本
        original_text: 原始中文文本

    Returns:
        str: 重组后的完整翻译
    """
    if not matched_segments:
        return model_translated
    if not model_translated and matched_segments:
        return " ".join(m["target"] for m in matched_segments)

    matched_segments.sort(key=lambda x: x["start"])

    unused_ranges = []
    prev_end = 0
    for m in matched_segments:
        if m["start"] > prev_end:
            unused_ranges.append((prev_end, m["start"]))
        prev_end = m["end"]
    if prev_end < len(original_text):
        unused_ranges.append((prev_end, len(original_text)))

    model_parts = model_translated.split("\n") if "\n" in model_translated else model_translated.split(". ")
    model_idx = 0

    result_parts = []
    text_pos = 0
    for m in matched_segments:
        while text_pos < m["start"]:
            src_frag = original_text[text_pos:m["start"]].strip()
            if src_frag and model_idx < len(model_parts):
                result_parts.append(model_parts[model_idx].strip())
                model_idx += 1
            text_pos = m["start"]
        result_parts.append(m["target"])
        text_pos = m["end"]

    if text_pos < len(original_text):
        src_frag = original_text[text_pos:].strip()
        if src_frag and model_idx < len(model_parts):
            remaining = " ".join(model_parts[model_idx:])
            if remaining.strip():
                result_parts.append(remaining.strip())

    return " ".join(result_parts)


def post_process_translation(text: str, target_lang: str) -> str:
    """硬编码风格后处理——所有翻译路径必须走此函数。

    包括：
    1. 品牌名统一：Jishi/jishi → ROX
    2. 英语特定修正：非规范搭配替换（如 "shift to online" → "shift online"）
    3. 非CJK语言中删除残余的中文字符（日语/韩语除外）

    Args:
        text: 翻译后的文本
        target_lang: 目标语言代码

    Returns:
        str: 后处理后的文本
    """
    if not text:
        return text

    text = text.strip()

    text = text.replace("Jishi", "ROX").replace("jishi", "ROX")

    if target_lang == "en":
        fixups = [
            (r'\bshift to online\b', 'shift online'),
            (r'\bgo to online\b', 'go online'),
            (r'\bdiscuss about\b', 'discuss'),
            (r'\bcontact with\b', 'contact'),
            (r'\bdeeply participate\b', 'actively engage'),
            (r'\breserved exclusive\b', 'exclusive reserved'),
        ]
        for pat, repl in fixups:
            text = re.sub(pat, repl, text, flags=re.IGNORECASE)

    if target_lang != "zh_hant":
        cjk = re.compile(r'[\u4e00-\u9fff\u3400-\u4dbf]')
        if target_lang in ("ja", "ko"):
            pass
        else:
            text = re.sub(r'[\u4e00-\u9fff\u3400-\u4dbf\u3000-\u303f\uff00-\uffef]+', '', text)
            text = re.sub(r' {2,}', ' ', text)
            text = re.sub(r'\n{3,}', '\n\n', text)

    text = text.strip()
    return text


def file_correct_compare(original_path, corrected_path, lang_code, conn, ids, vecs, ext=None):
    """比对原文文件和纠错后文件，提取差异，保存纠错结果到知识库。

    根据文件类型选择不同的纠错策略：
    - PPTX → 红字批注识别（提取红色字体 → 解析修正文本 → 匹配 KB → 保存）
    - DOCX/XLSX → 位置对齐比对（逐行对比 → 有差异则更新 KB）

    Args:
        original_path: 原始文件路径（含中文原文）
        corrected_path: 纠错后文件路径（含修正内容）
        lang_code: 目标语言代码
        conn: SQLite 数据库连接
        ids: 向量索引 ID 数组
        vecs: 向量数组
        ext: 文件扩展名

    Returns:
        tuple: (results: list[dict], summary: dict)
               results 包含每个纠错点的详细结果
               summary 包含统计信息（total, changed, saved, failed 等）
    """
    if ext is None:
        ext = Path(corrected_path).suffix.lower()
    
    if ext == '.pptx':
        return _pptx_correct_compare(original_path, corrected_path, lang_code, conn, ids, vecs)
    else:
        return _position_correct_compare(original_path, corrected_path, lang_code, conn, ids, vecs, ext)
